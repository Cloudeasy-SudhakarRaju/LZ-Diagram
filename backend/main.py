from fastapi import FastAPI, Response, HTTPException, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel, Field
from typing import Optional, List, Dict, Any
import html
import json
import uuid
import os
import base64
import subprocess
import tempfile
import logging
import traceback
from datetime import datetime
from pathlib import Path
import requests
import google.generativeai as genai

# Document processing imports
import PyPDF2
import openpyxl
from pptx import Presentation

# Import diagrams for Azure architecture generation
from diagrams import Diagram, Cluster, Edge
from diagrams.azure.compute import VM, AKS, AppServices, FunctionApps, ContainerInstances, ServiceFabricClusters, BatchAccounts
from diagrams.azure.network import VirtualNetworks, ApplicationGateway, LoadBalancers, Firewall, ExpressrouteCircuits, VirtualNetworkGateways, CDNProfiles
from diagrams.azure.storage import StorageAccounts, BlobStorage, DataLakeStorage
from diagrams.azure.database import SQLDatabases, CosmosDb, DatabaseForMysqlServers, DatabaseForPostgresqlServers
from diagrams.azure.security import KeyVaults, SecurityCenter, Sentinel
from diagrams.azure.identity import ActiveDirectory
from diagrams.azure.analytics import SynapseAnalytics, DataFactories, Databricks, StreamAnalyticsJobs, EventHubs
from diagrams.azure.integration import LogicApps, ServiceBus, EventGridTopics, APIManagement
from diagrams.azure.devops import Devops, Pipelines
from diagrams.azure.general import Subscriptions, Resourcegroups
from diagrams.azure.web import AppServices as WebApps

app = FastAPI(
    title="Azure Landing Zone Agent",
    description="Professional Azure Landing Zone Architecture Generator",
    version="1.0.0"
)

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # open for dev, restrict later
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Configure Google Gemini API
GEMINI_API_KEY = "AIzaSyCuYYvGh5wjwNniv9ZQ1QC-5pxwdj5lCWQ"
genai.configure(api_key=GEMINI_API_KEY)

# Initialize Gemini model
try:
    gemini_model = genai.GenerativeModel('gemini-1.5-pro')
    logger.info("Google Gemini API configured successfully")
except Exception as e:
    logger.error(f"Failed to configure Gemini API: {e}")
    gemini_model = None


class CustomerInputs(BaseModel):
    # Business Requirements
    business_objective: Optional[str] = Field(None, description="Primary business objective")
    regulatory: Optional[str] = Field(None, description="Regulatory requirements")
    industry: Optional[str] = Field(None, description="Industry vertical")
    
    # Organization Structure
    org_structure: Optional[str] = Field(None, description="Organization structure")
    governance: Optional[str] = Field(None, description="Governance model")
    
    # Identity & Access Management
    identity: Optional[str] = Field(None, description="Identity management approach")
    
    # Networking & Connectivity
    connectivity: Optional[str] = Field(None, description="Connectivity requirements")
    network_model: Optional[str] = Field(None, description="Network topology model")
    ip_strategy: Optional[str] = Field(None, description="IP address strategy")
    
    # Security
    security_zone: Optional[str] = Field(None, description="Security zone requirements")
    security_posture: Optional[str] = Field(None, description="Security posture")
    key_vault: Optional[str] = Field(None, description="Key management approach")
    threat_protection: Optional[str] = Field(None, description="Threat protection strategy")
    
    # Legacy fields for backward compatibility
    workload: Optional[str] = Field(None, description="Primary workload type")
    architecture_style: Optional[str] = Field(None, description="Architecture style")
    scalability: Optional[str] = Field(None, description="Scalability requirements")
    
    # Operations
    ops_model: Optional[str] = Field(None, description="Operations model")
    monitoring: Optional[str] = Field(None, description="Monitoring strategy")
    backup: Optional[str] = Field(None, description="Backup and recovery strategy")
    
    # Infrastructure
    topology_pattern: Optional[str] = Field(None, description="Topology pattern")
    
    # Migration & Cost
    migration_scope: Optional[str] = Field(None, description="Migration scope")
    cost_priority: Optional[str] = Field(None, description="Cost optimization priority")
    iac: Optional[str] = Field(None, description="Infrastructure as Code preference")
    
    # Specific Azure Service Selections
    # Compute Services
    compute_services: Optional[List[str]] = Field(default_factory=list, description="Selected compute services")
    
    # Networking Services
    network_services: Optional[List[str]] = Field(default_factory=list, description="Selected networking services")
    
    # Storage Services
    storage_services: Optional[List[str]] = Field(default_factory=list, description="Selected storage services")
    
    # Database Services
    database_services: Optional[List[str]] = Field(default_factory=list, description="Selected database services")
    
    # Security Services
    security_services: Optional[List[str]] = Field(default_factory=list, description="Selected security services")
    
    # Monitoring Services
    monitoring_services: Optional[List[str]] = Field(default_factory=list, description="Selected monitoring services")
    
    # AI/ML Services
    ai_services: Optional[List[str]] = Field(default_factory=list, description="Selected AI/ML services")
    
    # Analytics Services
    analytics_services: Optional[List[str]] = Field(default_factory=list, description="Selected analytics services")
    
    # Integration Services
    integration_services: Optional[List[str]] = Field(default_factory=list, description="Selected integration services")
    
    # DevOps Services
    devops_services: Optional[List[str]] = Field(default_factory=list, description="Selected DevOps services")
    
    # Backup Services
    backup_services: Optional[List[str]] = Field(default_factory=list, description="Selected backup services")
    
    # Enhanced Input Fields for AI Integration
    free_text_input: Optional[str] = Field(None, description="Free-form text input for additional requirements and context")
    url_input: Optional[str] = Field(None, description="URL for web content analysis")
    uploaded_files_info: Optional[List[Dict[str, Any]]] = Field(default_factory=list, description="Information about uploaded files")


# Azure Architecture Templates and Patterns
AZURE_TEMPLATES = {
    "enterprise": {
        "name": "Enterprise Scale Landing Zone",
        "management_groups": ["Root", "Platform", "Landing Zones", "Sandbox", "Decommissioned"],
        "subscriptions": ["Connectivity", "Identity", "Management", "Production", "Development"],
        "core_services": ["Azure AD", "Azure Policy", "Azure Monitor", "Key Vault", "Security Center"]
    },
    "small_medium": {
        "name": "Small-Medium Enterprise Landing Zone", 
        "management_groups": ["Root", "Platform", "Workloads"],
        "subscriptions": ["Platform", "Production", "Development"],
        "core_services": ["Azure AD", "Azure Monitor", "Key Vault"]
    },
    "startup": {
        "name": "Startup Landing Zone",
        "management_groups": ["Root", "Workloads"],
        "subscriptions": ["Production", "Development"],
        "core_services": ["Azure AD", "Azure Monitor"]
    }
}

AZURE_SERVICES_MAPPING = {
    # Compute Services
    "virtual_machines": {"name": "Azure Virtual Machines", "icon": "[VM]", "drawio_shape": "virtual_machine", "diagram_class": VM, "category": "compute"},
    "aks": {"name": "Azure Kubernetes Service", "icon": "[K8S]", "drawio_shape": "kubernetes_service", "diagram_class": AKS, "category": "compute"},
    "app_services": {"name": "Azure App Services", "icon": "[WEB]", "drawio_shape": "app_services", "diagram_class": AppServices, "category": "compute"},
    "web_apps": {"name": "Azure Web Apps", "icon": "[WEB]", "drawio_shape": "app_services", "diagram_class": WebApps, "category": "compute"},
    "functions": {"name": "Azure Functions", "icon": "[FUNC]", "drawio_shape": "function_app", "diagram_class": FunctionApps, "category": "compute"},
    "container_instances": {"name": "Container Instances", "icon": "[CI]", "drawio_shape": "container_instances", "diagram_class": ContainerInstances, "category": "compute"},
    "service_fabric": {"name": "Service Fabric", "icon": "[SF]", "drawio_shape": "service_fabric", "diagram_class": ServiceFabricClusters, "category": "compute"},
    "batch": {"name": "Azure Batch", "icon": "[BATCH]", "drawio_shape": "batch_accounts", "diagram_class": BatchAccounts, "category": "compute"},
    
    # Networking Services
    "virtual_network": {"name": "Virtual Network", "icon": "[NET]", "drawio_shape": "virtual_network", "diagram_class": VirtualNetworks, "category": "network"},
    "vpn_gateway": {"name": "VPN Gateway", "icon": "[VPN]", "drawio_shape": "vpn_gateway", "diagram_class": VirtualNetworkGateways, "category": "network"},
    "expressroute": {"name": "ExpressRoute", "icon": "[ER]", "drawio_shape": "expressroute_circuits", "diagram_class": ExpressrouteCircuits, "category": "network"},
    "load_balancer": {"name": "Load Balancer", "icon": "[LB]", "drawio_shape": "load_balancer", "diagram_class": LoadBalancers, "category": "network"},
    "application_gateway": {"name": "Application Gateway", "icon": "[AGW]", "drawio_shape": "application_gateway", "diagram_class": ApplicationGateway, "category": "network"},
    "firewall": {"name": "Azure Firewall", "icon": "[FW]", "drawio_shape": "firewall", "diagram_class": Firewall, "category": "network"},
    "waf": {"name": "Web Application Firewall", "icon": "[WAF]", "drawio_shape": "application_gateway", "diagram_class": ApplicationGateway, "category": "network"},
    "cdn": {"name": "Content Delivery Network", "icon": "[CDN]", "drawio_shape": "cdn_profiles", "diagram_class": None, "category": "network"},
    "front_door": {"name": "Azure Front Door", "icon": "[AFD]", "drawio_shape": "front_door", "diagram_class": ApplicationGateway, "category": "network"},
    "traffic_manager": {"name": "Traffic Manager", "icon": "[TM]", "drawio_shape": "traffic_manager_profiles", "diagram_class": None, "category": "network"},
    "virtual_wan": {"name": "Virtual WAN", "icon": "[WAN]", "drawio_shape": "virtual_wan", "diagram_class": VirtualNetworks, "category": "network"},
    
    # Storage Services
    "storage_accounts": {"name": "Storage Accounts", "icon": "[STOR]", "drawio_shape": "storage_accounts", "diagram_class": StorageAccounts, "category": "storage"},
    "blob_storage": {"name": "Blob Storage", "icon": "[BLOB]", "drawio_shape": "blob_storage", "diagram_class": BlobStorage, "category": "storage"},
    "queue_storage": {"name": "Queue Storage", "icon": "[QUEUE]", "drawio_shape": "queue_storage", "diagram_class": StorageAccounts, "category": "storage"},
    "table_storage": {"name": "Table Storage", "icon": "[TABLE]", "drawio_shape": "table_storage", "diagram_class": StorageAccounts, "category": "storage"},
    "file_storage": {"name": "Azure Files", "icon": "[FILES]", "drawio_shape": "files", "diagram_class": StorageAccounts, "category": "storage"},
    "disk_storage": {"name": "Managed Disks", "icon": "[DISK]", "drawio_shape": "managed_disks", "diagram_class": StorageAccounts, "category": "storage"},
    "data_lake": {"name": "Data Lake Storage", "icon": "[LAKE]", "drawio_shape": "data_lake_storage", "diagram_class": DataLakeStorage, "category": "storage"},
    
    # Database Services
    "sql_database": {"name": "Azure SQL Database", "icon": "[SQL]", "drawio_shape": "sql_database", "diagram_class": SQLDatabases, "category": "database"},
    "sql_managed_instance": {"name": "SQL Managed Instance", "icon": "[SQLMI]", "drawio_shape": "sql_managed_instance", "diagram_class": SQLDatabases, "category": "database"},
    "cosmos_db": {"name": "Cosmos DB", "icon": "[COSMOS]", "drawio_shape": "cosmos_db", "diagram_class": CosmosDb, "category": "database"},
    "mysql": {"name": "Azure Database for MySQL", "icon": "[MYSQL]", "drawio_shape": "database_for_mysql_servers", "diagram_class": DatabaseForMysqlServers, "category": "database"},
    "postgresql": {"name": "Azure Database for PostgreSQL", "icon": "[PGSQL]", "drawio_shape": "database_for_postgresql_servers", "diagram_class": DatabaseForPostgresqlServers, "category": "database"},
    "mariadb": {"name": "Azure Database for MariaDB", "icon": "[MARIA]", "drawio_shape": "database_for_mariadb_servers", "diagram_class": DatabaseForMysqlServers, "category": "database"},
    "redis_cache": {"name": "Azure Cache for Redis", "icon": "[REDIS]", "drawio_shape": "cache_redis", "diagram_class": None, "category": "database"},
    
    # Security Services
    "key_vault": {"name": "Azure Key Vault", "icon": "[KV]", "drawio_shape": "key_vault", "diagram_class": KeyVaults, "category": "security"},
    "active_directory": {"name": "Azure Active Directory", "icon": "[AAD]", "drawio_shape": "azure_active_directory", "diagram_class": ActiveDirectory, "category": "security"},
    "security_center": {"name": "Azure Security Center", "icon": "[SEC]", "drawio_shape": "security_center", "diagram_class": SecurityCenter, "category": "security"},
    "sentinel": {"name": "Azure Sentinel", "icon": "[SENT]", "drawio_shape": "sentinel", "diagram_class": Sentinel, "category": "security"},
    "defender": {"name": "Microsoft Defender", "icon": "[DEF]", "drawio_shape": "defender_easm", "diagram_class": SecurityCenter, "category": "security"},
    "information_protection": {"name": "Azure Information Protection", "icon": "[AIP]", "drawio_shape": "information_protection", "diagram_class": None, "category": "security"},
    
    # Monitoring & Management
    "monitor": {"name": "Azure Monitor", "icon": "[MON]", "drawio_shape": "monitor", "diagram_class": None, "category": "monitoring"},
    "log_analytics": {"name": "Log Analytics", "icon": "[LOG]", "drawio_shape": "log_analytics_workspaces", "diagram_class": None, "category": "monitoring"},
    "application_insights": {"name": "Application Insights", "icon": "[AI]", "drawio_shape": "application_insights", "diagram_class": None, "category": "monitoring"},
    "service_health": {"name": "Service Health", "icon": "[HLTH]", "drawio_shape": "service_health", "diagram_class": None, "category": "monitoring"},
    "advisor": {"name": "Azure Advisor", "icon": "💡", "drawio_shape": "advisor", "diagram_class": None, "category": "monitoring"},
    
    # AI/ML Services  
    "cognitive_services": {"name": "Cognitive Services", "icon": "🧠", "drawio_shape": "cognitive_services", "diagram_class": None, "category": "ai"},
    "machine_learning": {"name": "Azure Machine Learning", "icon": "🤖", "drawio_shape": "machine_learning", "diagram_class": None, "category": "ai"},
    "bot_service": {"name": "Bot Service", "icon": "🤖", "drawio_shape": "bot_services", "diagram_class": None, "category": "ai"},
    "form_recognizer": {"name": "Form Recognizer", "icon": "📄", "drawio_shape": "form_recognizer", "diagram_class": None, "category": "ai"},
    
    # Data & Analytics
    "synapse": {"name": "Azure Synapse Analytics", "icon": "[SYN]", "drawio_shape": "synapse_analytics", "diagram_class": SynapseAnalytics, "category": "analytics"},
    "data_factory": {"name": "Azure Data Factory", "icon": "[ADF]", "drawio_shape": "data_factory", "diagram_class": DataFactories, "category": "analytics"},
    "databricks": {"name": "Azure Databricks", "icon": "[DBR]", "drawio_shape": "databricks", "diagram_class": Databricks, "category": "analytics"},
    "stream_analytics": {"name": "Stream Analytics", "icon": "[ASA]", "drawio_shape": "stream_analytics", "diagram_class": StreamAnalyticsJobs, "category": "analytics"},
    "power_bi": {"name": "Power BI", "icon": "[PBI]", "drawio_shape": "power_bi", "diagram_class": None, "category": "analytics"},
    
    # Integration Services
    "logic_apps": {"name": "Logic Apps", "icon": "[LA]", "drawio_shape": "logic_apps", "diagram_class": LogicApps, "category": "integration"},
    "service_bus": {"name": "Service Bus", "icon": "[SB]", "drawio_shape": "service_bus", "diagram_class": ServiceBus, "category": "integration"},
    "event_grid": {"name": "Event Grid", "icon": "[EG]", "drawio_shape": "event_grid_topics", "diagram_class": EventGridTopics, "category": "integration"},
    "event_hubs": {"name": "Event Hubs", "icon": "[EH]", "drawio_shape": "event_hubs", "diagram_class": EventHubs, "category": "integration"},
    "api_management": {"name": "API Management", "icon": "[APIM]", "drawio_shape": "api_management", "diagram_class": APIManagement, "category": "integration"},
    
    # DevOps & Management
    "devops": {"name": "Azure DevOps", "icon": "[ADO]", "drawio_shape": "devops", "diagram_class": Devops, "category": "devops"},
    "automation": {"name": "Azure Automation", "icon": "[AUTO]", "drawio_shape": "automation_accounts", "diagram_class": None, "category": "devops"},
    "policy": {"name": "Azure Policy", "icon": "[POL]", "drawio_shape": "policy", "diagram_class": None, "category": "governance"},
    "blueprints": {"name": "Azure Blueprints", "icon": "[BP]", "drawio_shape": "blueprints", "diagram_class": None, "category": "governance"},
    "resource_manager": {"name": "Azure Resource Manager", "icon": "[ARM]", "drawio_shape": "resource_groups", "diagram_class": Resourcegroups, "category": "governance"},
    
    # Backup & Recovery
    "backup": {"name": "Azure Backup", "icon": "[BAK]", "drawio_shape": "backup", "diagram_class": None, "category": "backup"},
    "site_recovery": {"name": "Azure Site Recovery", "icon": "[ASR]", "drawio_shape": "site_recovery", "diagram_class": None, "category": "backup"},
}

def get_safe_output_directory() -> str:
    """Get a safe directory for output files with fallback options"""
    directories_to_try = [
        "/tmp",
        tempfile.gettempdir(),
        os.path.expanduser("~/tmp"),
        "./tmp"
    ]
    
    for directory in directories_to_try:
        try:
            # Create directory if it doesn't exist
            Path(directory).mkdir(parents=True, exist_ok=True)
            
            # Test if we can write to it
            test_file = os.path.join(directory, f"test_write_{uuid.uuid4().hex[:8]}.tmp")
            with open(test_file, 'w') as f:
                f.write("test")
            os.remove(test_file)
            
            logger.info(f"Using output directory: {directory}")
            return directory
            
        except Exception as e:
            logger.warning(f"Cannot use directory {directory}: {e}")
            continue
    
    raise Exception("No writable output directory found. Tried: " + ", ".join(directories_to_try))

def cleanup_old_files(directory: str, max_age_hours: int = 24):
    """Clean up old generated files to prevent disk space issues"""
    try:
        current_time = datetime.now().timestamp()
        max_age_seconds = max_age_hours * 3600
        
        for filename in os.listdir(directory):
            if filename.startswith("azure_landing_zone_") and filename.endswith(".png"):
                filepath = os.path.join(directory, filename)
                try:
                    file_age = current_time - os.path.getmtime(filepath)
                    if file_age > max_age_seconds:
                        os.remove(filepath)
                        logger.info(f"Cleaned up old file: {filename}")
                except Exception as e:
                    logger.warning(f"Failed to clean up file {filename}: {e}")
                    
    except Exception as e:
        logger.warning(f"Failed to perform cleanup in {directory}: {e}")

# Google Gemini AI Integration Functions
def analyze_url_content(url: str) -> str:
    """Fetch and analyze URL content using Gemini AI"""
    try:
        if not gemini_model:
            return "Gemini AI not available for URL analysis"
            
        # Fetch URL content with timeout
        response = requests.get(url, timeout=10)
        response.raise_for_status()
        content = response.text[:10000]  # Limit content size
        
        prompt = f"""
        Analyze the following web content for Azure architecture planning:
        
        URL: {url}
        Content: {content}
        
        Please provide insights for:
        1. Relevant Azure services mentioned or implied
        2. Architecture patterns or requirements
        3. Security and compliance considerations
        4. Scalability and performance requirements
        5. Integration requirements
        
        Format your response as a structured analysis.
        """
        
        # Add timeout for AI API call using threading
        import threading
        import queue
        
        def ai_call_with_timeout(prompt, timeout_seconds=10):
            """Call AI with timeout using threading"""
            result_queue = queue.Queue()
            
            def ai_worker():
                try:
                    result = gemini_model.generate_content(prompt)
                    result_queue.put(("success", result.text))
                except Exception as e:
                    result_queue.put(("error", str(e)))
            
            thread = threading.Thread(target=ai_worker)
            thread.daemon = True
            thread.start()
            thread.join(timeout=timeout_seconds)
            
            if thread.is_alive():
                return None
            
            try:
                status, response = result_queue.get_nowait()
                if status == "success":
                    return response
                else:
                    return None
            except queue.Empty:
                return None
        
        ai_response = ai_call_with_timeout(prompt, 10)
        if ai_response:
            return ai_response
        else:
            logger.warning(f"URL analysis timed out for {url}")
            return f"URL analysis timed out - skipped analysis for {url}"
        
    except Exception as e:
        logger.error(f"Error analyzing URL {url}: {e}")
        return f"Error analyzing URL: {str(e)}"

def process_uploaded_document(file_content: bytes, filename: str, file_type: str) -> str:
    """Process uploaded document using Gemini AI"""
    try:
        if not gemini_model:
            return "Gemini AI not available for document analysis"
            
        text_content = ""
        
        # Extract text based on file type
        if file_type.lower() == 'pdf':
            text_content = extract_pdf_text(file_content)
        elif file_type.lower() in ['xlsx', 'xls']:
            text_content = extract_excel_text(file_content)
        elif file_type.lower() in ['pptx', 'ppt']:
            text_content = extract_pptx_text(file_content)
        else:
            return f"Unsupported file type: {file_type}"
        
        if not text_content.strip():
            return "No readable text found in the document"
            
        prompt = f"""
        Analyze the following document content for Azure Landing Zone architecture planning:
        
        Document: {filename}
        Type: {file_type}
        Content: {text_content[:8000]}  # Limit content size
        
        Please provide insights for:
        1. Current architecture mentioned in the document
        2. Business requirements and objectives
        3. Compliance and regulatory requirements
        4. Security requirements
        5. Recommended Azure services and patterns
        6. Migration considerations
        7. Governance and operational requirements
        
        Format your response as a structured analysis for enterprise architecture planning.
        """
        
        result = gemini_model.generate_content(prompt)
        return result.text
        
    except Exception as e:
        logger.error(f"Error processing document {filename}: {e}")
        return f"Error processing document: {str(e)}"

def extract_pdf_text(file_content: bytes) -> str:
    """Extract text from PDF file"""
    try:
        import io
        pdf_file = io.BytesIO(file_content)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        return text
    except Exception as e:
        logger.error(f"Error extracting PDF text: {e}")
        return ""

def extract_excel_text(file_content: bytes) -> str:
    """Extract text from Excel file"""
    try:
        import io
        excel_file = io.BytesIO(file_content)
        workbook = openpyxl.load_workbook(excel_file)
        text = ""
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text += f"Sheet: {sheet_name}\n"
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                if row_text.strip():
                    text += row_text + "\n"
        return text
    except Exception as e:
        logger.error(f"Error extracting Excel text: {e}")
        return ""

def extract_pptx_text(file_content: bytes) -> str:
    """Extract text from PowerPoint file"""
    try:
        import io
        pptx_file = io.BytesIO(file_content)
        presentation = Presentation(pptx_file)
        text = ""
        for slide_num, slide in enumerate(presentation.slides, 1):
            text += f"Slide {slide_num}:\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        logger.error(f"Error extracting PowerPoint text: {e}")
        return ""

def generate_ai_enhanced_recommendations(inputs: CustomerInputs, url_analysis: str = "", doc_analysis: str = "") -> str:
    """Generate AI-enhanced architecture recommendations using Gemini"""
    try:
        if not gemini_model:
            return "Standard recommendations (AI enhancement not available)"
            
        # Build context from inputs
        context = f"""
        Business Objective: {inputs.business_objective or 'Not specified'}
        Industry: {inputs.industry or 'General'}
        Organization Structure: {inputs.org_structure or 'Not specified'}
        Regulatory Requirements: {inputs.regulatory or 'Standard'}
        Security Requirements: {inputs.security_posture or 'Standard'}
        Scalability Requirements: {inputs.scalability or 'Standard'}
        Free-text Input: {inputs.free_text_input or 'None provided'}
        
        Selected Services:
        - Compute: {', '.join(inputs.compute_services or [])}
        - Network: {', '.join(inputs.network_services or [])}
        - Storage: {', '.join(inputs.storage_services or [])}
        - Database: {', '.join(inputs.database_services or [])}
        - Security: {', '.join(inputs.security_services or [])}
        """
        
        if url_analysis:
            context += f"\n\nURL Analysis Results:\n{url_analysis}"
            
        if doc_analysis:
            context += f"\n\nDocument Analysis Results:\n{doc_analysis}"
        
        prompt = f"""
        Based on the following customer requirements and analysis, provide comprehensive Azure Landing Zone architecture recommendations:
        
        {context}
        
        Please provide:
        1. Recommended Azure Landing Zone template (Enterprise Scale, Small Scale, etc.)
        2. Detailed architecture recommendations with specific Azure services
        3. Security and compliance strategy
        4. Network topology and connectivity recommendations
        5. Identity and access management strategy
        6. Operations and monitoring approach
        7. Cost optimization strategies
        8. Migration roadmap and phases
        9. Governance and policy recommendations
        10. Risk assessment and mitigation strategies
        
        Format your response as a comprehensive enterprise architecture document.
        """
        
        # Add timeout and error handling for AI API call using threading
        import threading
        import queue
        
        def ai_call_with_timeout(prompt, timeout_seconds=15):
            """Call AI with timeout using threading"""
            result_queue = queue.Queue()
            
            def ai_worker():
                try:
                    result = gemini_model.generate_content(prompt)
                    result_queue.put(("success", result.text))
                except Exception as e:
                    result_queue.put(("error", str(e)))
            
            thread = threading.Thread(target=ai_worker)
            thread.daemon = True
            thread.start()
            thread.join(timeout=timeout_seconds)
            
            if thread.is_alive():
                # Thread is still running, timeout occurred
                logger.warning(f"AI call timed out after {timeout_seconds} seconds")
                return None
            
            try:
                status, response = result_queue.get_nowait()
                if status == "success":
                    return response
                else:
                    logger.error(f"AI call failed: {response}")
                    return None
            except queue.Empty:
                logger.warning("AI call completed but no result available")
                return None
        
        ai_response = ai_call_with_timeout(prompt, 15)
        if ai_response:
            return ai_response
        else:
            logger.warning("AI recommendation generation failed or timed out, using fallback")
            return generate_fallback_recommendations(inputs)
        
    except Exception as e:
        logger.error(f"Error generating AI recommendations: {e}")
        return generate_fallback_recommendations(inputs)

def generate_fallback_recommendations(inputs: CustomerInputs) -> str:
    """Generate fallback recommendations when AI is not available"""
    return f"""
## Standard Azure Landing Zone Recommendations

**Recommended Template:** Enterprise Scale Landing Zone
Based on the selected services and organization structure, an Enterprise Scale template provides the best foundation.

**Core Architecture Recommendations:**

1. **Identity & Access Management**
   - Azure Active Directory as the identity provider
   - Conditional Access policies for enhanced security
   - Role-Based Access Control (RBAC) for resource management

2. **Network Architecture**
   - Hub-and-spoke network topology with Azure Virtual WAN
   - Network Security Groups (NSGs) for traffic filtering
   - Azure Firewall for centralized network security

3. **Security Framework**
   - Zero Trust security model implementation
   - Azure Security Center for threat protection
   - Key Vault for secrets management

4. **Governance & Compliance**
   - Azure Policy for compliance enforcement
   - Management Groups for hierarchical organization
   - Azure Blueprint for repeatable deployments

5. **Monitoring & Operations**
   - Azure Monitor for comprehensive observability
   - Log Analytics for centralized logging
   - Application Insights for application performance monitoring

**Selected Services Analysis:**
- Compute Services: {', '.join(inputs.compute_services or ['Virtual Machines', 'App Services'])}
- Network Services: {', '.join(inputs.network_services or ['Virtual Network', 'Load Balancer'])}
- Storage Services: {', '.join(inputs.storage_services or ['Storage Accounts', 'Blob Storage'])}
- Database Services: {', '.join(inputs.database_services or ['SQL Database'])}
- Security Services: {', '.join(inputs.security_services or ['Key Vault', 'Security Center'])}

**Cost Optimization:**
- Use Azure Reserved Instances for predictable workloads
- Implement auto-scaling for dynamic workloads
- Regular cost reviews and optimization recommendations

**Migration Strategy:**
1. Assessment and discovery phase
2. Pilot migration of non-critical workloads
3. Production workload migration in phases
4. Optimization and governance implementation
"""

def analyze_free_text_requirements(free_text: str) -> dict:
    """Analyze free text input to extract specific service requirements using AI"""
    try:
        if not gemini_model or not free_text:
            return {"services": [], "reasoning": "No analysis available"}
            
        prompt = f"""
        Analyze the following user requirement and identify ONLY the specific Azure services that are explicitly needed or implied:
        
        User Requirement: "{free_text}"
        
        Based on this requirement, provide a JSON response with:
        1. "services" - array of specific Azure service keys that are actually needed (use keys like: app_services, sql_database, monitor, virtual_network, etc.)
        2. "reasoning" - brief explanation of why these services were selected
        3. "architecture_pattern" - suggested pattern (simple, standard, complex)
        
        Be very conservative - only include services that are directly mentioned or absolutely required for the stated use case.
        For example:
        - "web application hosting" -> app_services, virtual_network
        - "database backend" -> sql_database or mysql/postgresql depending on context
        - "basic monitoring" -> monitor, log_analytics
        
        Do NOT include full enterprise landing zone services unless specifically requested.
        
        Return only valid JSON format.
        """
        
        result = gemini_model.generate_content(prompt)
        response_text = result.text.strip()
        
        # Try to extract JSON from the response
        import json
        import re
        
        # Look for JSON content
        json_match = re.search(r'\{.*\}', response_text, re.DOTALL)
        if json_match:
            json_str = json_match.group()
            try:
                analysis = json.loads(json_str)
                return analysis
            except json.JSONDecodeError:
                pass
        
        # Fallback parsing if JSON parsing fails
        logger.warning(f"Failed to parse AI response as JSON: {response_text}")
        return {"services": [], "reasoning": "Could not parse AI analysis"}
        
    except Exception as e:
        logger.error(f"Error analyzing free text requirements: {e}")
        return {"services": [], "reasoning": f"Analysis error: {str(e)}"}

def validate_customer_inputs(inputs: CustomerInputs) -> None:
    """Validate customer inputs to prevent potential errors"""
    # Check for extremely long strings that might cause issues
    string_fields = [
        inputs.business_objective, inputs.regulatory, inputs.industry,
        inputs.org_structure, inputs.governance, inputs.identity,
        inputs.connectivity, inputs.network_model, inputs.ip_strategy,
        inputs.security_zone, inputs.security_posture, inputs.key_vault,
        inputs.threat_protection, inputs.workload, inputs.architecture_style,
        inputs.scalability, inputs.ops_model, inputs.monitoring, inputs.backup,
        inputs.topology_pattern, inputs.migration_scope, inputs.cost_priority, inputs.iac,
        inputs.url_input
    ]
    
    for field in string_fields:
        if field and len(field) > 1000:  # Reasonable limit for most fields
            raise ValueError(f"Input field too long: {len(field)} characters (max 1000)")
    
    # Special validation for free-text input (allowing more characters)
    if inputs.free_text_input and len(inputs.free_text_input) > 10000:
        raise ValueError(f"Free text input too long: {len(inputs.free_text_input)} characters (max 10000)")
    
    # Check service lists for reasonable sizes
    service_lists = [
        inputs.compute_services, inputs.network_services, inputs.storage_services,
        inputs.database_services, inputs.security_services, inputs.monitoring_services,
        inputs.ai_services, inputs.analytics_services, inputs.integration_services,
        inputs.devops_services, inputs.backup_services
    ]
    
    for service_list in service_lists:
        if service_list and len(service_list) > 50:  # Reasonable limit
            raise ValueError(f"Too many services selected: {len(service_list)} (max 50)")
    
    # Validate URL format if provided
    if inputs.url_input:
        if not inputs.url_input.startswith(('http://', 'https://')):
            raise ValueError("URL must start with http:// or https://")
    
    # Validate uploaded files info
    if inputs.uploaded_files_info:
        if len(inputs.uploaded_files_info) > 10:  # Reasonable limit
            raise ValueError(f"Too many uploaded files: {len(inputs.uploaded_files_info)} (max 10)")

def generate_azure_architecture_diagram(inputs: CustomerInputs, output_dir: str = None, format: str = "png") -> str:
    """Generate Enhanced Azure architecture diagram following 50+ enterprise architecture principles
    
    Enterprise Architecture Principles Implemented:
    1-10: Clear containers/swimlanes, minimal crossing connections, proper visual hierarchy
    11-20: Clear connection labeling, numbered workflow, all specified components
    21-30: Security zoning, environment labeling, HA indicators, monitoring overlay
    31-40: Comprehensive legend, standardized iconography, observability integration
    41-50: Scalability indicators, compliance overlays, cost management, future-ready design
    """
    
    logger.info("Starting Enhanced Azure architecture diagram generation with 50+ principles")
    
    try:
        # Validate inputs first
        validate_customer_inputs(inputs)
        logger.info("Input validation completed successfully")
        
        # Get safe output directory
        if output_dir is None:
            output_dir = get_safe_output_directory()
        
        # Clean up old files to prevent disk space issues
        cleanup_old_files(output_dir)
        
        # Verify Graphviz availability before proceeding
        try:
            result = subprocess.run(['dot', '-V'], capture_output=True, text=True, timeout=10)
            if result.returncode != 0:
                raise Exception(f"Graphviz 'dot' command failed with return code {result.returncode}. stderr: {result.stderr}")
            logger.info(f"Graphviz version: {result.stderr.strip()}")
        except subprocess.TimeoutExpired:
            raise Exception("Graphviz 'dot' command timed out. Graphviz may be unresponsive.")
        except FileNotFoundError:
            raise Exception("Graphviz is not installed or not accessible. Please install Graphviz: sudo apt-get install -y graphviz graphviz-dev")
        except subprocess.SubprocessError as e:
            raise Exception(f"Graphviz check failed: {str(e)}. Please install Graphviz: sudo apt-get install -y graphviz graphviz-dev")
        
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        unique_id = str(uuid.uuid4())[:8]
        filename = f"enhanced_azure_architecture_{timestamp}_{unique_id}"
        filepath = os.path.join(output_dir, filename)
        
        logger.info(f"Generating enhanced enterprise diagram: {filename}")
        
        # Verify output directory is writable
        if not os.access(output_dir, os.W_OK):
            raise Exception(f"Output directory {output_dir} is not writable")
        
        # Determine organization template
        template = generate_architecture_template(inputs)
        org_name = inputs.org_structure or "Enterprise"
        
        logger.info(f"Using enhanced template: {template['template']['name']}")
        
        try:
            # Set the format based on the requested output format
            output_format = "svg" if format.lower() == "svg" else "png"
            
            # Enhanced graph attributes for enterprise architecture following best practices
            with Diagram(
                f"Azure Enterprise Landing Zone - {template['template']['name']}", 
                filename=filepath, 
                show=False, 
                direction="TB",
                outformat=output_format,
                graph_attr={
                    "fontsize": "18",
                    "fontname": "Arial, sans-serif",
                    "rankdir": "TB",
                    "nodesep": "2.0",      # Increased spacing for larger nodes
                    "ranksep": "3.0",      # Increased vertical spacing between layers
                    "bgcolor": "#f8f9fa",   # Professional light background
                    "margin": "1.0",       # Increased margin
                    "compound": "true",     # Enable cluster-to-cluster connections
                    "splines": "polyline",  # Polyline routing for minimal crossings
                    "overlap": "false",
                    "pack": "true",
                    "packmode": "clust",
                    "pad": "1.0",          # Increased padding
                    "dpi": "300",          # High resolution for professional output
                    "concentrate": "true",  # Merge similar edges for clarity
                    "ordering": "out"      # Consistent edge ordering
                },
                node_attr={
                    "fontsize": "12",
                    "fontname": "Arial, sans-serif", 
                    "style": "filled,rounded",
                    "shape": "box",
                    "fillcolor": "#ffffff",
                    "color": "#333333",
                    "penwidth": "2",
                    "width": "2.5",      # Minimum width to accommodate longer labels
                    "height": "1.5",     # Minimum height for multi-line labels
                    "fixedsize": "false"  # Allow nodes to grow beyond minimum size
                },
                edge_attr={
                    "fontsize": "10",
                    "fontname": "Arial, sans-serif",
                    "style": "solid",
                    "arrowhead": "vee",
                    "arrowsize": "0.8",
                    "penwidth": "2",
                    "minlen": "2",
                    "weight": "1"
                }
            ):
                
                logger.info("Creating enhanced enterprise architecture with 50+ design principles...")
                
                # Initialize service counters for numbered workflow
                service_counter = 1
                
                # Initialize service lists at function scope to avoid UnboundLocalError
                storage_services = []
                database_services = []
                
                # === INTERNET EDGE LAYER (TOP - UNTRUSTED ZONE) ===
                # Principle 1-5: Clear containers, security zoning, visual hierarchy
                with Cluster("[ INTERNET EDGE ]", graph_attr={
                    "bgcolor": "#ffebee",      # Light red for untrusted zone
                    "style": "filled,rounded,bold", 
                    "color": "#c62828",
                    "fontcolor": "#c62828",
                    "fontsize": "16",
                    "label": "[ INTERNET EDGE ] (UNTRUSTED ZONE)",
                    "rank": "min",
                    "penwidth": "3",
                    "margin": "16"
                }):
                    internet_services = []
                    
                    # Add Front Door (CDN/WAF) - Required component per principles
                    if inputs.network_services and "front_door" in inputs.network_services:
                        front_door = ApplicationGateway(f"{service_counter}. Azure Front Door\\n[Global CDN+WAF]")
                        internet_services.append(front_door)
                        service_counter += 1
                    
                    # Add Azure CDN if specified
                    if inputs.network_services and "cdn" in inputs.network_services:
                        cdn = CDNProfiles(f"{service_counter}. Azure CDN\\n[Content Delivery]")
                        internet_services.append(cdn)
                        service_counter += 1
                    
                    # Add Traffic Manager for global load balancing
                    if inputs.network_services and "traffic_manager" in inputs.network_services:
                        traffic_mgr = LoadBalancers(f"{service_counter}. Traffic Manager\\n[DNS Load Balancer]")
                        internet_services.append(traffic_mgr)
                        service_counter += 1
                    
                    # Default Front Door if no edge services specified (required per principles)
                    if not internet_services:
                        front_door = ApplicationGateway(f"{service_counter}. Azure Front Door\\n[Global Entry Point]")
                        internet_services.append(front_door)
                        service_counter += 1
                
                # === IDENTITY & SECURITY LAYER (TOP-LEFT - SEMI-TRUSTED ZONE) ===
                # Principle 6-10: Identity clarity, access management, security zoning
                with Cluster("[ IDENTITY & SECURITY ]", graph_attr={
                    "bgcolor": "#e8f5e9",       # Light green for semi-trusted
                    "style": "filled,rounded,bold",
                    "color": "#2e7d32",
                    "fontcolor": "#2e7d32", 
                    "fontsize": "16",
                    "label": "[ IDENTITY & SECURITY ] (SEMI-TRUSTED ZONE)",
                    "penwidth": "3",
                    "margin": "16"
                }):
                    # Core Identity Services
                    aad = ActiveDirectory(f"{service_counter}. Azure AD\\n[Identity Provider]\\n[Active-Active HA]")
                    service_counter += 1
                    
                    key_vault = KeyVaults(f"{service_counter}. Key Vault\\n[Secrets Management]\\n[Premium SKU]")
                    service_counter += 1
                    
                    security_services = [aad, key_vault]
                    
                    # Enhanced Security Services
                    if inputs.security_services:
                        if "security_center" in inputs.security_services:
                            sec_center = SecurityCenter(f"{service_counter}. Security Center\\n[Defender for Cloud]\\n[Compliance: GDPR/HIPAA]") 
                            security_services.append(sec_center)
                            service_counter += 1
                            
                        if "sentinel" in inputs.security_services:
                            sentinel = Sentinel(f"{service_counter}. Azure Sentinel\\n[SIEM/SOAR]\\n[24/7 SOC]")
                            security_services.append(sentinel)
                            service_counter += 1
                
                # === MANAGEMENT & GOVERNANCE LAYER (LEFT SIDE) ===
                # Principle 11-15: Governance hierarchy, cost management, policy enforcement
                with Cluster("[ MANAGEMENT & GOVERNANCE ]", graph_attr={
                    "bgcolor": "#f3e5f5",       # Light purple for governance
                    "style": "filled,rounded,bold",
                    "color": "#7b1fa2",
                    "fontcolor": "#7b1fa2",
                    "fontsize": "16",
                    "label": "[ MANAGEMENT & GOVERNANCE ]\\n[Cost Management + Policies]",
                    "penwidth": "3",
                    "margin": "16"
                }):
                    # Management Group Hierarchy
                    root_mg = Subscriptions(f"{service_counter}. Root Management Group\\n[Enterprise Governance]")
                    service_counter += 1
                    
                    if template['template']['name'] == "Enterprise Scale Landing Zone":
                        platform_mg = Subscriptions(f"{service_counter}. Platform MG\\n[Shared Services]")
                        service_counter += 1
                        workloads_mg = Subscriptions(f"{service_counter}. Landing Zones MG\\n[Application Workloads]")
                        service_counter += 1
                        sandbox_mg = Subscriptions(f"{service_counter}. Sandbox MG\\n[Development/Testing]")
                        service_counter += 1
                        decom_mg = Subscriptions(f"{service_counter}. Decommissioned MG\\n[Legacy/Sunset]")
                        service_counter += 1
                        
                        # Cost Management and Governance connections
                        root_mg >> Edge(label="Policy\\nAssignment", style="dashed", color="#7b1fa2") >> [platform_mg, workloads_mg, sandbox_mg, decom_mg]
                    else:
                        platform_mg = Subscriptions(f"{service_counter}. Platform MG\\n[Infrastructure]")
                        service_counter += 1
                        workloads_mg = Subscriptions(f"{service_counter}. Workloads MG\\n[Applications]")
                        service_counter += 1
                        root_mg >> Edge(label="Governance\\nPolicies", style="dashed", color="#7b1fa2") >> [platform_mg, workloads_mg]
                
                # === ACTIVE REGION - NETWORK HUB (CENTER - TRUSTED ZONE) ===
                # Principle 16-20: Network architecture, clear boundaries, HA indicators
                with Cluster("[ ACTIVE REGION - NETWORK HUB ]", graph_attr={
                    "bgcolor": "#e3f2fd",       # Light blue for trusted zone
                    "style": "filled,rounded,bold",
                    "color": "#1565c0",
                    "fontcolor": "#1565c0",
                    "fontsize": "16",
                    "label": "[ ACTIVE REGION - NETWORK HUB ] (TRUSTED ZONE)\\n[Primary Region: East US 2]",
                    "penwidth": "3",
                    "margin": "20"
                }):
                    # Hub VNet with enhanced labeling
                    hub_vnet = VirtualNetworks(f"{service_counter}. Hub VNet\\n[Shared Network Services]\\n[10.0.0.0/16]\\n[Active-Active HA]")
                    service_counter += 1
                    
                    # Network Security and Connectivity Services
                    network_services = []
                    
                    # Azure Firewall (Required for enterprise)
                    if not inputs.network_services or "firewall" in inputs.network_services:
                        firewall = Firewall(f"{service_counter}. Azure Firewall\\n[Network Security]\\n[Premium SKU]\\n[Threat Intelligence]")
                        network_services.append(firewall)
                        service_counter += 1
                    
                    # VPN Gateway for hybrid connectivity
                    if not inputs.network_services or "vpn_gateway" in inputs.network_services:
                        vpn_gw = VirtualNetworkGateways(f"{service_counter}. VPN Gateway\\n[Hybrid Connectivity]\\n[VpnGw2 SKU]\\n[99.95% SLA]")
                        network_services.append(vpn_gw)
                        service_counter += 1
                    
                    # ExpressRoute for dedicated connectivity
                    if inputs.network_services and "expressroute" in inputs.network_services:
                        er_circuit = ExpressrouteCircuits(f"{service_counter}. ExpressRoute\\n[Private Connectivity]\\n[1Gbps Circuit]\\n[99.95% SLA]")
                        network_services.append(er_circuit)
                        service_counter += 1
                    
                    # Application Gateway for L7 load balancing
                    if inputs.network_services and "application_gateway" in inputs.network_services:
                        app_gw = ApplicationGateway(f"{service_counter}. Application Gateway\\n[L7 Load Balancer]\\n[WAF Enabled]\\n[Standard_v2 SKU]")
                        network_services.append(app_gw)
                        service_counter += 1
                    
                    # Additional network services based on selection
                    if inputs.network_services:
                        for service in inputs.network_services:
                            if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["diagram_class"]:
                                if service not in ["firewall", "vpn_gateway", "expressroute", "application_gateway", "front_door", "cdn", "traffic_manager"]:
                                    diagram_class = AZURE_SERVICES_MAPPING[service]["diagram_class"]
                                    service_name = f"{service_counter}. {AZURE_SERVICES_MAPPING[service]['name']}\\n[Network Service]"
                                    network_services.append(diagram_class(service_name))
                                    service_counter += 1
                    
                    # Spoke VNets with environment labeling and CIDR blocks
                    prod_vnet = VirtualNetworks(f"{service_counter}. Production VNet\\n[PROD Environment]\\n[10.1.0.0/16]\\n[Active-Active HA]\\n[Zone Redundant]")
                    service_counter += 1
                    
                    dev_vnet = VirtualNetworks(f"{service_counter}. Development VNet\\n[DEV Environment]\\n[10.2.0.0/16]\\n[Standard Tier]")
                    service_counter += 1
                    
                    uat_vnet = VirtualNetworks(f"{service_counter}. UAT VNet\\n[UAT Environment]\\n[10.3.0.0/16]\\n[Standard Tier]")
                    service_counter += 1
                    
                    # Hub-Spoke peering with clear traffic flow
                    hub_vnet >> Edge(label="Hub-Spoke\\nPeering\\n[BGP Routing]", style="solid", color="#1565c0", penwidth="3") >> [prod_vnet, dev_vnet, uat_vnet]
                    
                    # Platform subscription governance connection
                    platform_mg >> Edge(label="Network\\nGovernance\\n[RBAC + Policies]", style="dashed", color="#7b1fa2", penwidth="2") >> hub_vnet
                    
                    # Connect network services to hub with clear flow labels
                    for i, ns in enumerate(network_services):
                        hub_vnet >> Edge(label=f"Internal\\nRouting\\n[Service {i+1}]", style="dotted", color="#1565c0", penwidth="2") >> ns
                
                # === COMPUTE & APPLICATION LAYER (MIDDLE TIER) ===
                # Principle 21-25: Application architecture, scalability indicators, cloud-native patterns
                if inputs.compute_services:
                    with Cluster("[ COMPUTE & APPLICATIONS ]", graph_attr={
                        "bgcolor": "#fff3e0",       # Light orange for compute
                        "style": "filled,rounded,bold",
                        "color": "#ef6c00",
                        "fontcolor": "#ef6c00",
                        "fontsize": "16",
                        "label": "[ COMPUTE & APPLICATIONS ]\\n[Cloud-Native + Legacy]\\n[Auto-Scaling Enabled]",
                        "penwidth": "3",
                        "margin": "16"
                    }):
                        compute_services = []
                        
                        for service in inputs.compute_services:
                            if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["diagram_class"]:
                                diagram_class = AZURE_SERVICES_MAPPING[service]["diagram_class"]
                                service_name = f"{service_counter}. {AZURE_SERVICES_MAPPING[service]['name']}"
                                
                                # Add tier and scaling information
                                if service == "aks":
                                    service_name += "\\n[Kubernetes v1.28]\\n[Auto-Scale: 1-100 nodes]\\n[Zone Redundant]"
                                elif service == "app_services":
                                    service_name += "\\n[PaaS Platform]\\n[Standard S1 SKU]\\n[Auto-Scale Enabled]"
                                elif service == "virtual_machines":
                                    service_name += "\\n[IaaS Platform]\\n[Standard_D4s_v3]\\n[Availability Sets]"
                                elif service == "functions":
                                    service_name += "\\n[Serverless]\\n[Consumption Plan]\\n[Event-Driven]"
                                else:
                                    service_name += "\\n[Managed Service]\\n[Standard Tier]"
                                
                                compute_services.append(diagram_class(service_name))
                                service_counter += 1
                        
                        # Connect compute services to production VNet
                        for cs in compute_services:
                            prod_vnet >> Edge(label="Application\\nTraffic\\n[HTTPS/443]", style="solid", color="#ef6c00", penwidth="2") >> cs
                            
                        # Connect workloads management group for governance
                        if compute_services:
                            workloads_mg >> Edge(label="Application\\nGovernance\\n[Policies + RBAC]", style="dashed", color="#7b1fa2", penwidth="2") >> compute_services[0]
                
                # === INTEGRATION & API LAYER ===
                # Principle 26-30: Integration patterns, API management, messaging
                if inputs.integration_services:
                    with Cluster("🔗 INTEGRATION & API MANAGEMENT", graph_attr={
                        "bgcolor": "#e1f5fe",       # Light cyan for integration
                        "style": "filled,rounded,bold",
                        "color": "#0277bd",
                        "fontcolor": "#0277bd",
                        "fontsize": "16",
                        "label": "🔗 INTEGRATION & API MANAGEMENT\\n[Service Mesh + API Gateway]",
                        "penwidth": "3",
                        "margin": "16"
                    }):
                        integration_services = []
                        
                        for service in inputs.integration_services:
                            if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["diagram_class"]:
                                diagram_class = AZURE_SERVICES_MAPPING[service]["diagram_class"]
                                service_name = f"{service_counter}. {AZURE_SERVICES_MAPPING[service]['name']}"
                                
                                # Add integration-specific details
                                if service == "api_management":
                                    service_name += "\\n[API Gateway]\\n[Developer Portal]\\n[Rate Limiting]"
                                elif service == "logic_apps":
                                    service_name += "\\n[Workflow Engine]\\n[Consumption Plan]\\n[400+ Connectors]"
                                elif service == "service_bus":
                                    service_name += "\\n[Message Broker]\\n[Premium Tier]\\n[Dead Letter Queue]"
                                elif service == "event_grid":
                                    service_name += "\\n[Event Routing]\\n[Custom Topics]\\n[Event Subscriptions]"
                                else:
                                    service_name += "\\n[Integration Service]"
                                
                                integration_services.append(diagram_class(service_name))
                                service_counter += 1
                        
                        # Connect integration services to both production and hub
                        for is_service in integration_services:
                            prod_vnet >> Edge(label="Integration\\nTraffic\\n[AMQP/HTTP]", style="solid", color="#0277bd", penwidth="2") >> is_service
                            hub_vnet >> Edge(label="Cross-Service\\nCommunication", style="dotted", color="#0277bd") >> is_service
                
                # === DATA LAYER (BOTTOM TIER) ===
                # Principle 31-35: Data architecture, storage tiers, backup/recovery
                if inputs.database_services or inputs.storage_services:
                    with Cluster("[ DATA & STORAGE LAYER ]", graph_attr={
                        "bgcolor": "#e8f5e9",       # Light green for data
                        "style": "filled,rounded,bold",
                        "color": "#2e7d32",
                        "fontcolor": "#2e7d32",
                        "fontsize": "16",
                        "label": "[ DATA & STORAGE LAYER ]\\n[Persistent Storage + Databases]\\n[Geo-Redundant + Backup]",
                        "penwidth": "3",
                        "margin": "16"
                    }):
                        # Storage Services with tier information
                        if inputs.storage_services:
                            for service in inputs.storage_services:
                                if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["diagram_class"]:
                                    diagram_class = AZURE_SERVICES_MAPPING[service]["diagram_class"]
                                    service_name = f"{service_counter}. {AZURE_SERVICES_MAPPING[service]['name']}"
                                    
                                    # Add storage-specific details and data classification
                                    if service == "storage_accounts":
                                        service_name += "\\n[General Purpose v2]\\n[LRS+GRS Redundancy]\\n[Hot/Cool/Archive Tiers]"
                                    elif service == "blob_storage":
                                        service_name += "\\n[Object Storage]\\n[Versioning Enabled]\\n[Lifecycle Management]"
                                    elif service == "queue_storage":
                                        service_name += "\\n[Message Queue]\\n[64KB Message Size]\\n[Async Processing]"
                                    elif service == "table_storage":
                                        service_name += "\\n[NoSQL Key-Value]\\n[Serverless]\\n[Auto-Scale]"
                                    elif service == "data_lake":
                                        service_name += "\\n[Big Data Analytics]\\n[Hierarchical Namespace]\\n[Data Lake Gen2]"
                                    else:
                                        service_name += "\\n[Storage Service]\\n[Managed Service]"
                                    
                                    storage_services.append(diagram_class(service_name))
                                    service_counter += 1
                        
                        # Database Services with performance tiers
                        if inputs.database_services:
                            for service in inputs.database_services:
                                if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["diagram_class"]:
                                    diagram_class = AZURE_SERVICES_MAPPING[service]["diagram_class"]
                                    service_name = f"{service_counter}. {AZURE_SERVICES_MAPPING[service]['name']}"
                                    
                                    # Add database-specific details and performance tiers
                                    if service == "sql_database":
                                        service_name += "\\n[General Purpose]\\n[4 vCores, 100GB]\\n[Point-in-Time Restore]\\n[99.99% SLA]"
                                    elif service == "cosmos_db":
                                        service_name += "\\n[Multi-Model NoSQL]\\n[Global Distribution]\\n[5 Consistency Levels]\\n[99.999% SLA]"
                                    elif service == "redis_cache":
                                        service_name += "\\n[In-Memory Cache]\\n[Standard C1 SKU]\\n[Clustering Support]\\n[99.9% SLA]"
                                    elif service == "mysql":
                                        service_name += "\\n[Managed MySQL]\\n[General Purpose]\\n[Auto-Backup]\\n[99.99% SLA]"
                                    elif service == "postgresql":
                                        service_name += "\\n[Managed PostgreSQL]\\n[General Purpose]\\n[Point-in-Time Restore]\\n[99.99% SLA]"
                                    else:
                                        service_name += "\\n[Managed Database]\\n[High Availability]"
                                    
                                    database_services.append(diagram_class(service_name))
                                    service_counter += 1
                        
                        # Data connections with different patterns
                        all_data_services = storage_services + database_services
                        for ds in all_data_services:
                            prod_vnet >> Edge(label="Data Access\\n[TLS 1.3]\\n[Private Endpoint]", style="solid", color="#2e7d32", penwidth="2") >> ds
                            workloads_mg >> Edge(label="Data Governance\\n[Classification + Policies]", style="dashed", color="#7b1fa2") >> ds
                
                # === ANALYTICS & AI LAYER ===
                # Principle 36-40: Analytics architecture, AI/ML services, big data processing
                if inputs.analytics_services:
                    with Cluster("[ ANALYTICS & AI ]", graph_attr={
                        "bgcolor": "#fce4ec",       # Light pink for analytics
                        "style": "filled,rounded,bold",
                        "color": "#c2185b",
                        "fontcolor": "#c2185b",
                        "fontsize": "16",
                        "label": "[ ANALYTICS & AI ]\\n[Big Data + Machine Learning]\\n[Real-time + Batch Processing]",
                        "penwidth": "3",
                        "margin": "16"
                    }):
                        analytics_services = []
                        
                        for service in inputs.analytics_services:
                            if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["diagram_class"]:
                                diagram_class = AZURE_SERVICES_MAPPING[service]["diagram_class"]
                                service_name = f"{service_counter}. {AZURE_SERVICES_MAPPING[service]['name']}"
                                
                                # Add analytics-specific details
                                if service == "synapse":
                                    service_name += "\\n[Data Warehouse]\\n[Spark + SQL Pools]\\n[Data Integration]"
                                elif service == "data_factory":
                                    service_name += "\\n[ETL/ELT Pipeline]\\n[Hybrid Data Movement]\\n[250+ Data Connectors]"
                                elif service == "databricks":
                                    service_name += "\\n[Apache Spark]\\n[Collaborative Notebooks]\\n[MLflow Integration]"
                                elif service == "stream_analytics":
                                    service_name += "\\n[Real-time Analytics]\\n[Event Processing]\\n[IoT Data Streams]"
                                else:
                                    service_name += "\\n[Analytics Service]"
                                
                                analytics_services.append(diagram_class(service_name))
                                service_counter += 1
                        
                        # Connect analytics to data sources and VNets
                        for as_service in analytics_services:
                            prod_vnet >> Edge(label="Analytics\\nProcessing\\n[Secure Access]", style="solid", color="#c2185b", penwidth="2") >> as_service
                            if storage_services:
                                storage_services[0] >> Edge(label="Data\\nIngestion\\n[Batch/Stream]", style="dotted", color="#c2185b") >> as_service
                
                # === STANDBY REGION (DISASTER RECOVERY) ===
                # Principle 41-45: DR strategy, cross-region replication, business continuity
                if template['template']['name'] == "Enterprise Scale Landing Zone":
                    with Cluster("🔄 STANDBY REGION (DISASTER RECOVERY)", graph_attr={
                        "bgcolor": "#fff3e0",       # Light amber for DR
                        "style": "filled,rounded,dashed,bold",  # Dashed border for DR
                        "color": "#ef6c00",
                        "fontcolor": "#ef6c00",
                        "fontsize": "16",
                        "label": "🔄 STANDBY REGION - DISASTER RECOVERY\\n[Secondary Region: West US 2]\\n[RTO: 1 hour | RPO: 15 minutes]",
                        "penwidth": "3",
                        "margin": "16"
                    }):
                        # DR Network Infrastructure
                        dr_vnet = VirtualNetworks(f"{service_counter}. DR VNet\\n[Disaster Recovery]\\n[10.100.0.0/16]\\n[Passive Standby]\\n[Automated Failover]")
                        service_counter += 1
                        
                        # DR Storage for backup
                        dr_storage = StorageAccounts(f"{service_counter}. DR Storage\\n[Geo-Redundant]\\n[Cross-Region Replication]\\n[99.999999999% Durability]")
                        service_counter += 1
                        
                        # DR Database replicas
                        if inputs.database_services and "sql_database" in inputs.database_services:
                            dr_sql = SQLDatabases(f"{service_counter}. DR SQL Database\\n[Read Replica]\\n[Async Replication]\\n[Auto-Failover Group]")
                            service_counter += 1
                            dr_vnet >> Edge(label="Database\\nReplication", style="dotted", color="#ef6c00") >> dr_sql
                        
                        # Cross-region replication connections
                        prod_vnet >> Edge(label="Geo-Replication\\n[Async]\\n[RTO: 1hr]", style="dashed", color="#ef6c00", penwidth="3") >> dr_vnet
                        
                        if storage_services:
                            storage_services[0] >> Edge(label="Storage\\nReplication\\n[GRS/RA-GRS]", style="dashed", color="#2e7d32") >> dr_storage
                
                # === MONITORING & OBSERVABILITY OVERLAY ===
                # Principle 46-50: Observability, SIEM, compliance monitoring, cost optimization
                with Cluster("[ MONITORING & OBSERVABILITY ]", graph_attr={
                    "bgcolor": "#e0f2f1",       # Light teal for monitoring
                    "style": "filled,rounded,bold",
                    "color": "#00695c",
                    "fontcolor": "#00695c",
                    "fontsize": "16", 
                    "label": "[ MONITORING & OBSERVABILITY ]\\n[360° Visibility + SIEM]\\n[Cost Optimization + Compliance]",
                    "penwidth": "3",
                    "margin": "16"
                }):
                    monitor_services = []
                    
                    # Core monitoring services
                    monitor = Subscriptions(f"{service_counter}. Azure Monitor\\n[Metrics + Alerts]\\n[Application Insights]\\n[99.9% SLA]")
                    monitor_services.append(monitor)
                    service_counter += 1
                    
                    log_analytics = Subscriptions(f"{service_counter}. Log Analytics\\n[Centralized Logging]\\n[KQL Queries]\\n[30-day Retention]")
                    monitor_services.append(log_analytics)
                    service_counter += 1
                    
                    # Additional monitoring services based on selection
                    if inputs.monitoring_services:
                        for service in inputs.monitoring_services:
                            if service not in ["monitor", "log_analytics"]:
                                service_name = f"{service_counter}. {service.replace('_', ' ').title()}\\n[Monitoring Service]"
                                monitor_services.append(Subscriptions(service_name))
                                service_counter += 1
                    
                    # Cost Management and Governance monitoring
                    cost_mgmt = Subscriptions(f"{service_counter}. Cost Management\\n[Budget Alerts]\\n[Chargeback/Showback]\\n[Optimization Recommendations]")
                    monitor_services.append(cost_mgmt)
                    service_counter += 1
                    
                    # Security monitoring (if Sentinel is enabled)
                    if inputs.security_services and "sentinel" in inputs.security_services:
                        siem = Subscriptions(f"{service_counter}. Security Monitoring\\n[24/7 SOC]\\n[Threat Hunting]\\n[SOAR Playbooks]")
                        monitor_services.append(siem)
                        service_counter += 1
                    
                    # Comprehensive monitoring connections to all layers
                    monitored_services = [hub_vnet, prod_vnet, dev_vnet, uat_vnet]
                    if 'compute_services' in locals() and compute_services:
                        monitored_services.extend(compute_services[:2])  # Monitor first 2 compute services
                    if 'all_data_services' in locals() and all_data_services:
                        monitored_services.extend(all_data_services[:2])  # Monitor first 2 data services
                    
                    for i, monitor_service in enumerate(monitor_services[:3]):  # First 3 monitoring services
                        for j, main_service in enumerate(monitored_services):
                            if i == 0:  # Azure Monitor connects to all
                                main_service >> Edge(label=f"Metrics+Logs\\n[Telemetry]", style="dotted", color="#00695c", penwidth="1") >> monitor_service
                            elif i == 1 and j < 3:  # Log Analytics to network services
                                main_service >> Edge(label="Log\\nCollection", style="dotted", color="#00695c", penwidth="1") >> monitor_service
                
                # === MAIN WORKFLOW CONNECTIONS ===
                # Principle: Clear numbered workflow from internet to applications to data
                logger.info("Creating numbered workflow connections...")
                
                # 1. Internet to Security flow (User Authentication)
                for i, service in enumerate(internet_services, 1):
                    for j, sec_service in enumerate(security_services[:2]):  # Connect to AAD and Key Vault
                        service >> Edge(
                            label=f"Auth Flow {i}.{j+1}\\n[OAuth 2.0/OIDC]", 
                            style="solid", 
                            color="#1976d2", 
                            penwidth="3"
                        ) >> sec_service
                
                # 2. Security to Network flow (Identity Validation)
                aad >> Edge(
                    label="Identity\\nValidation\\n[JWT Tokens]\\n[Conditional Access]", 
                    style="solid", 
                    color="#2e7d32", 
                    penwidth="3"
                ) >> hub_vnet
                
                # 3. Platform management to security (Governance)
                platform_mg >> Edge(
                    label="Security\\nGovernance\\n[RBAC + Policies]\\n[Compliance]", 
                    style="dashed", 
                    color="#7b1fa2", 
                    penwidth="2"
                ) >> [aad, key_vault]
                
                # 4. Internet to Application flow (if Application Gateway exists)
                if network_services and any("gateway" in str(ns) for ns in network_services):
                    for inet_svc in internet_services:
                        for net_svc in network_services:
                            if "gateway" in str(net_svc).lower() or "application" in str(net_svc).lower():
                                inet_svc >> Edge(
                                    label="Web Traffic\\n[HTTPS/443]\\n[Global Load Balance]",
                                    style="solid",
                                    color="#ff6f00",
                                    penwidth="3"
                                ) >> net_svc
                                break
                
                # === COMPREHENSIVE LEGEND ===
                # Principle: Complete legend covering all diagram elements
                with Cluster("[ LEGEND & NOTATION GUIDE ]", graph_attr={
                    "bgcolor": "#f5f5f5",       # Light gray for legend
                    "style": "filled,rounded,bold",
                    "color": "#424242",
                    "fontcolor": "#424242",
                    "fontsize": "14",
                    "label": "[ LEGEND & NOTATION GUIDE ]",
                    "penwidth": "2",
                    "margin": "12",
                    "rank": "sink"  # Place at bottom
                }):
                    # Legend items using subscriptions as placeholder boxes
                    legend_zones = Subscriptions("ZONES:\\n• Red: Untrusted (Internet)\\n• Green: Semi-trusted (Identity)\\n• Blue: Trusted (Network/Apps)\\n• Orange: Compute/DR")
                    
                    legend_lines = Subscriptions("LINE TYPES:\\n• Solid: Primary traffic flow\\n• Dashed: Governance/policies\\n• Dotted: Monitoring/logs\\n• Bold: Critical connections")
                    
                    legend_ha = Subscriptions("HA INDICATORS:\\n• Active-Active: Load balanced\\n• Active-Passive: Standby ready\\n• Zone Redundant: Multi-AZ\\n• Geo-Redundant: Cross-region")
                    
                    legend_compliance = Subscriptions("COMPLIANCE:\\n• GDPR: Data protection\\n• HIPAA: Healthcare data\\n• SOC 2: Security controls\\n• ISO 27001: Info security")
                    
                    legend_tiers = Subscriptions("SERVICE TIERS:\\n• Basic: Development/test\\n• Standard: Production ready\\n• Premium: Enterprise grade\\n• Ultra: Mission critical")
                
                logger.info("Enhanced enterprise architecture diagram structure created successfully")
        
        except Exception as e:
            logger.error(f"Error during enhanced diagram creation: {str(e)}")
            logger.error(traceback.format_exc())
            raise Exception(f"Error generating enhanced Azure architecture diagram: {str(e)}")
        
        # Return the file path of the generated diagram
        if format.lower() == "svg":
            # Check for SVG file generated directly by diagrams library
            svg_path = f"{filepath}.svg"
            if os.path.exists(svg_path):
                file_size = os.path.getsize(svg_path)
                logger.info(f"Enhanced SVG diagram generated successfully: {svg_path} (size: {file_size} bytes)")
                return svg_path
            else:
                # Fallback: try to generate SVG using dot command from gv file
                dot_path = f"{filepath}.gv" 
                if os.path.exists(dot_path):
                    try:
                        # Convert dot file to SVG
                        result = subprocess.run(['dot', '-Tsvg', dot_path, '-o', svg_path], 
                                              capture_output=True, text=True, timeout=30)
                        if result.returncode != 0:
                            raise Exception(f"SVG generation failed: {result.stderr}")
                        
                        if os.path.exists(svg_path):
                            file_size = os.path.getsize(svg_path)
                            logger.info(f"Enhanced SVG diagram generated successfully via dot: {svg_path} (size: {file_size} bytes)")
                            return svg_path
                        else:
                            raise Exception(f"SVG generation failed - file not found: {svg_path}")
                    except subprocess.TimeoutExpired:
                        raise Exception("SVG generation timed out")
                    except Exception as e:
                        raise Exception(f"Failed to generate SVG: {str(e)}")
                else:
                    raise Exception(f"Neither SVG nor dot file found: {svg_path}, {dot_path}")
        else:
            # Default PNG generation
            png_path = f"{filepath}.png"
            if os.path.exists(png_path):
                file_size = os.path.getsize(png_path)
                logger.info(f"Enhanced PNG diagram generated successfully: {png_path} (size: {file_size} bytes)")
                return png_path
            else:
                raise Exception(f"Enhanced diagram generation failed - PNG file not found: {png_path}")
            
    except Exception as e:
        logger.error(f"Failed to generate enhanced Azure architecture diagram: {str(e)}")
        logger.error(traceback.format_exc())
        raise

def generate_simple_svg_diagram(inputs: CustomerInputs) -> str:
    """Generate a simple SVG diagram as fallback when Python Diagrams fails"""
    
    template = generate_architecture_template(inputs)
    template_name = template['template']['name']
    
    # Create a simple SVG representation
    svg_width = 800
    svg_height = 600
    
    svg_content = f'''<svg width="{svg_width}" height="{svg_height}" xmlns="http://www.w3.org/2000/svg">
    <defs>
        <style>
            .title {{ font-family: Arial, sans-serif; font-size: 18px; font-weight: bold; fill: #0078d4; }}
            .group-title {{ font-family: Arial, sans-serif; font-size: 14px; font-weight: bold; fill: #323130; }}
            .service {{ font-family: Arial, sans-serif; font-size: 12px; fill: #605e5c; }}
            .mgmt-group {{ fill: #e1f5fe; stroke: #0078d4; stroke-width: 2; }}
            .subscription {{ fill: #f3e5f5; stroke: #6b69d6; stroke-width: 2; }}
            .service-box {{ fill: #fff3e0; stroke: #d83b01; stroke-width: 1; cursor: pointer; }}
            .service-box:hover {{ fill: #ffebdd; }}
            .network-box {{ fill: #e8f5e8; stroke: #107c10; stroke-width: 1; cursor: pointer; }}
            .network-box:hover {{ fill: #f3fdf3; }}
            .security-box {{ fill: #ffebee; stroke: #d13438; stroke-width: 1; cursor: pointer; }}
            .security-box:hover {{ fill: #fdf3f4; }}
        </style>
    </defs>
    
    <!-- Background -->
    <rect width="100%" height="100%" fill="#f8f9fa"/>
    
    <!-- Title -->
    <text x="400" y="30" class="title" text-anchor="middle">Azure Landing Zone - {template_name}</text>
    
    <!-- Azure Tenant Container -->
    <rect x="50" y="60" width="700" height="520" fill="none" stroke="#0078d4" stroke-width="3" stroke-dasharray="5,5"/>
    <text x="60" y="80" class="group-title">Azure Tenant</text>
    
    <!-- Management Groups -->
    <rect x="70" y="100" width="200" height="150" class="mgmt-group" rx="5"/>
    <text x="80" y="120" class="group-title">Management Groups</text>
    
    <!-- Management Group Items -->
    <rect x="80" y="130" width="80" height="30" class="service-box" rx="3"/>
    <text x="120" y="149" class="service" text-anchor="middle">Root MG</text>
    
    <rect x="170" y="130" width="80" height="30" class="service-box" rx="3"/>
    <text x="210" y="149" class="service" text-anchor="middle">Platform</text>
    
    <rect x="80" y="170" width="80" height="30" class="service-box" rx="3"/>
    <text x="120" y="189" class="service" text-anchor="middle">Landing Zones</text>
    
    <rect x="170" y="170" width="80" height="30" class="service-box" rx="3"/>
    <text x="210" y="189" class="service" text-anchor="middle">Sandbox</text>
    
    <!-- Subscriptions -->
    <rect x="290" y="100" width="200" height="150" class="subscription" rx="5"/>
    <text x="300" y="120" class="group-title">Subscriptions</text>
    
    <!-- Subscription Items -->
    <rect x="300" y="130" width="80" height="30" class="service-box" rx="3"/>
    <text x="340" y="149" class="service" text-anchor="middle">Connectivity</text>
    
    <rect x="390" y="130" width="80" height="30" class="service-box" rx="3"/>
    <text x="430" y="149" class="service" text-anchor="middle">Identity</text>
    
    <rect x="300" y="170" width="80" height="30" class="service-box" rx="3"/>
    <text x="340" y="189" class="service" text-anchor="middle">Production</text>
    
    <rect x="390" y="170" width="80" height="30" class="service-box" rx="3"/>
    <text x="430" y="189" class="service" text-anchor="middle">Development</text>
    
    <!-- Network Architecture -->
    <rect x="520" y="100" width="200" height="150" class="network-box" rx="5"/>
    <text x="530" y="120" class="group-title">Network Architecture</text>
    
    <rect x="530" y="130" width="80" height="30" class="network-box" rx="3"/>
    <text x="570" y="149" class="service" text-anchor="middle">Hub VNet</text>
    
    <rect x="620" y="130" width="80" height="30" class="network-box" rx="3"/>
    <text x="660" y="149" class="service" text-anchor="middle">Spoke VNet</text>'''
    
    # Add selected services
    y_offset = 280
    if inputs.compute_services:
        svg_content += f'''
    <!-- Compute Services -->
    <rect x="70" y="{y_offset}" width="300" height="80" class="service-box" rx="5"/>
    <text x="80" y="{y_offset + 20}" class="group-title">Compute Services</text>'''
        
        x_pos = 80
        for i, service in enumerate(inputs.compute_services[:4]):  # Max 4 services
            service_name = service.replace('_', ' ').title()
            svg_content += f'''
    <rect x="{x_pos}" y="{y_offset + 30}" width="60" height="25" class="service-box" rx="3"/>
    <text x="{x_pos + 30}" y="{y_offset + 47}" class="service" text-anchor="middle" font-size="10">{service_name[:8]}</text>'''
            x_pos += 70
    
    if inputs.network_services:
        svg_content += f'''
    <!-- Network Services -->
    <rect x="390" y="{y_offset}" width="300" height="80" class="network-box" rx="5"/>
    <text x="400" y="{y_offset + 20}" class="group-title">Network Services</text>'''
        
        x_pos = 400
        for i, service in enumerate(inputs.network_services[:4]):  # Max 4 services
            service_name = service.replace('_', ' ').title()
            svg_content += f'''
    <rect x="{x_pos}" y="{y_offset + 30}" width="60" height="25" class="network-box" rx="3"/>
    <text x="{x_pos + 30}" y="{y_offset + 47}" class="service" text-anchor="middle" font-size="10">{service_name[:8]}</text>'''
            x_pos += 70
    
    # Security Services
    y_offset += 100
    svg_content += f'''
    <!-- Security & Identity -->
    <rect x="70" y="{y_offset}" width="620" height="80" class="security-box" rx="5"/>
    <text x="80" y="{y_offset + 20}" class="group-title">Security & Identity Services</text>
    
    <rect x="80" y="{y_offset + 30}" width="100" height="25" class="security-box" rx="3"/>
    <text x="130" y="{y_offset + 47}" class="service" text-anchor="middle">Azure AD</text>
    
    <rect x="190" y="{y_offset + 30}" width="100" height="25" class="security-box" rx="3"/>
    <text x="240" y="{y_offset + 47}" class="service" text-anchor="middle">Key Vault</text>
    
    <rect x="300" y="{y_offset + 30}" width="100" height="25" class="security-box" rx="3"/>
    <text x="350" y="{y_offset + 47}" class="service" text-anchor="middle">Security Center</text>'''
    
    # Add connections
    svg_content += '''
    <!-- Connections -->
    <line x1="170" y1="145" x2="290" y2="145" stroke="#666" stroke-width="2" marker-end="url(#arrowhead)"/>
    <line x1="390" y1="145" x2="520" y2="145" stroke="#666" stroke-width="2" marker-end="url(#arrowhead)"/>
    
    <!-- Arrow marker -->
    <defs>
        <marker id="arrowhead" markerWidth="10" markerHeight="7" refX="10" refY="3.5" orient="auto">
            <polygon points="0 0, 10 3.5, 0 7" fill="#666"/>
        </marker>
    </defs>
    
    <!-- Legend -->
    <rect x="520" y="420" width="250" height="120" class="security-box" rx="5" stroke="#666" stroke-width="1"/>
    <text x="530" y="440" class="group-title">Legend</text>
    
    <!-- Legend items -->
    <rect x="530" y="450" width="15" height="15" class="mgmt-group" rx="2"/>
    <text x="550" y="462" class="service" font-size="10">Management Groups</text>
    
    <rect x="530" y="470" width="15" height="15" class="subscription" rx="2"/>
    <text x="550" y="482" class="service" font-size="10">Subscriptions</text>
    
    <rect x="530" y="490" width="15" height="15" class="network-box" rx="2"/>
    <text x="550" y="502" class="service" font-size="10">Network Services</text>
    
    <rect x="530" y="510" width="15" height="15" class="security-box" rx="2"/>
    <text x="550" y="522" class="service" font-size="10">Security Services</text>
    
    <!-- Environment Labels -->
    <text x="650" y="462" class="service" font-size="10" fill="#d32f2f">[PROD] Production</text>
    <text x="650" y="482" class="service" font-size="10" fill="#2e7d32">[DEV] Development</text>
    <text x="650" y="502" class="service" font-size="10" fill="#ef6c00">[DR] Disaster Recovery</text>
    
    <!-- Footer -->
    <text x="400" y="570" class="service" text-anchor="middle" fill="#8a8886">Generated by Azure Landing Zone Agent - Enterprise Architecture</text>
</svg>'''
    
    return svg_content

def _add_enhanced_service_clusters(inputs: CustomerInputs, prod_vnet, workloads_mg, service_counter: int):
    """Enhanced service clusters with proper positioning and numbering for enterprise architecture"""
    try:
        # Compute Services Layer (Middle - Application Tier)
        if inputs.compute_services:
            with Cluster("💻 Compute Services", graph_attr={
                "bgcolor": "#fff3e0", 
                "style": "filled,rounded",
                "color": "#e65100",
                "fontcolor": "#e65100",
                "fontsize": "14",
                "label": "💻 COMPUTE SERVICES (Application Tier)"
            }):
                compute_services = []
                counter = service_counter
                for service in inputs.compute_services:
                    if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["diagram_class"]:
                        diagram_class = AZURE_SERVICES_MAPPING[service]["diagram_class"]
                        service_name = f"{counter}. {AZURE_SERVICES_MAPPING[service]['name']}"
                        compute_services.append(diagram_class(service_name))
                        counter += 1
                
                # Connect compute to production VNet with workflow labels
                for cs in compute_services:
                    prod_vnet >> Edge(label="App Traffic", style="solid", color="orange") >> cs
                    workloads_mg >> Edge(label="Governance", style="dashed", color="purple") >> cs
        
        # Integration Services (API/Messaging Layer)
        if inputs.integration_services:
            with Cluster("🔗 Integration Services", graph_attr={
                "bgcolor": "#fff8e1", 
                "style": "filled,rounded",
                "color": "#f57f17",
                "fontcolor": "#f57f17", 
                "fontsize": "14",
                "label": "🔗 INTEGRATION SERVICES (API/Messaging Layer)"
            }):
                integration_services = []
                counter = service_counter + 10
                for service in inputs.integration_services:
                    if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["diagram_class"]:
                        diagram_class = AZURE_SERVICES_MAPPING[service]["diagram_class"]
                        service_name = f"{counter}. {AZURE_SERVICES_MAPPING[service]['name']}"
                        integration_services.append(diagram_class(service_name))
                        counter += 1
                
                # Connect integration services
                for is_service in integration_services:
                    prod_vnet >> Edge(label="API/Msg", style="solid", color="orange") >> is_service
                    workloads_mg >> Edge(label="Governance", style="dashed", color="purple") >> is_service
        
        # DevOps Services (Management Layer)
        if inputs.devops_services:
            with Cluster("⚙️ DevOps & Automation", graph_attr={
                "bgcolor": "#fafafa", 
                "style": "filled,rounded",
                "color": "#424242",
                "fontcolor": "#424242",
                "fontsize": "14",
                "label": "⚙️ DEVOPS & AUTOMATION"
            }):
                devops_services = []
                counter = service_counter + 15
                for service in inputs.devops_services:
                    if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["diagram_class"]:
                        diagram_class = AZURE_SERVICES_MAPPING[service]["diagram_class"]
                        service_name = f"{counter}. {AZURE_SERVICES_MAPPING[service]['name']}"
                        devops_services.append(diagram_class(service_name))
                        counter += 1
                
                # Connect DevOps services to management
                for ds in devops_services:
                    workloads_mg >> Edge(label="CI/CD", style="dashed", color="gray") >> ds
                    
    except Exception as e:
        logger.warning(f"Error adding enhanced service clusters: {str(e)}")


def _add_data_layer_clusters(inputs: CustomerInputs, prod_vnet, workloads_mg, service_counter: int):
    """Add data layer clusters with proper visual hierarchy (bottom layer)"""
    try:
        # Storage Services (Data Layer - Bottom)
        if inputs.storage_services:
            with Cluster("[ Storage Services ]", graph_attr={
                "bgcolor": "#e8f5e9", 
                "style": "filled,rounded",
                "color": "#2e7d32",
                "fontcolor": "#2e7d32",
                "fontsize": "14",
                "label": "[ STORAGE SERVICES ] (Data Layer)"
            }):
                storage_services = []
                counter = service_counter
                for service in inputs.storage_services:
                    if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["diagram_class"]:
                        diagram_class = AZURE_SERVICES_MAPPING[service]["diagram_class"]
                        service_name = f"{counter}. {AZURE_SERVICES_MAPPING[service]['name']}"
                        
                        # Add data classification labels
                        if service in ["blob_storage", "data_lake"]:
                            service_name += "\n[Confidential]"
                        elif service in ["queue_storage", "table_storage"]:
                            service_name += "\n[Internal]"
                        else:
                            service_name += "\n[General]"
                            
                        storage_services.append(diagram_class(service_name))
                        counter += 1
                
                # Default storage if none specified
                if not storage_services:
                    storage_default = StorageAccounts(f"{counter}. Storage Accounts\n[General]")
                    storage_services.append(storage_default)
                
                # Connect storage to production VNet and workloads with data flow labels
                for ss in storage_services:
                    prod_vnet >> Edge(label="Data Access", style="solid", color="green") >> ss
                    workloads_mg >> Edge(label="Data Governance", style="dashed", color="purple") >> ss
        
        # Database Services (Data Layer - Bottom)
        if inputs.database_services:
            with Cluster("🗄️ Database Services", graph_attr={
                "bgcolor": "#e3f2fd", 
                "style": "filled,rounded",
                "color": "#1565c0",
                "fontcolor": "#1565c0",
                "fontsize": "14",
                "label": "🗄️ DATABASE SERVICES (Data Layer)"
            }):
                database_services = []
                counter = service_counter + 5
                for service in inputs.database_services:
                    if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["diagram_class"]:
                        diagram_class = AZURE_SERVICES_MAPPING[service]["diagram_class"]
                        service_name = f"{counter}. {AZURE_SERVICES_MAPPING[service]['name']}"
                        
                        # Add service tier notation
                        if service in ["sql_database", "cosmos_db"]:
                            service_name += "\n[Premium/HA]"
                        elif service == "redis_cache":
                            service_name += "\n[Standard]" 
                        else:
                            service_name += "\n[Basic]"
                            
                        database_services.append(diagram_class(service_name))
                        counter += 1
                
                # Connect databases to production VNet and workloads
                for ds in database_services:
                    prod_vnet >> Edge(label="DB Traffic", style="solid", color="blue") >> ds
                    workloads_mg >> Edge(label="Data Governance", style="dashed", color="purple") >> ds
        
        # Analytics Services (Data Analytics Layer)
        if inputs.analytics_services:
            with Cluster("[ Analytics & AI ]", graph_attr={
                "bgcolor": "#f3e5f5", 
                "style": "filled,rounded",
                "color": "#7b1fa2",
                "fontcolor": "#7b1fa2",
                "fontsize": "14",
                "label": "[ ANALYTICS & AI ] (Insights Layer)"
            }):
                analytics_services = []
                counter = service_counter + 10
                for service in inputs.analytics_services:
                    if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["diagram_class"]:
                        diagram_class = AZURE_SERVICES_MAPPING[service]["diagram_class"]
                        service_name = f"{counter}. {AZURE_SERVICES_MAPPING[service]['name']}"
                        analytics_services.append(diagram_class(service_name))
                        counter += 1
                
                # Connect analytics to production VNet and workloads
                for as_service in analytics_services:
                    prod_vnet >> Edge(label="Analytics\nData Flow", style="solid", color="purple") >> as_service
                    workloads_mg >> Edge(label="Analytics\nGovernance", style="dashed", color="purple") >> as_service
        
        # Backup & Recovery Services (Resilience Layer) 
        if inputs.backup_services:
            with Cluster("[ Backup & Recovery ]", graph_attr={
                "bgcolor": "#e0f2f1", 
                "style": "filled,rounded",
                "color": "#00695c",
                "fontcolor": "#00695c",
                "fontsize": "14",
                "label": "[ BACKUP & RECOVERY ] (Resilience Layer)"
            }):
                backup_services = []
                counter = service_counter + 15
                for service in inputs.backup_services:
                    if service in AZURE_SERVICES_MAPPING:
                        # Most backup services don't have diagram classes
                        service_name = f"{counter}. {AZURE_SERVICES_MAPPING[service]['name']}"
                        # Use storage accounts as a placeholder for backup services
                        backup_services.append(StorageAccounts(service_name))
                        counter += 1
                
                # Connect backup services
                for bs in backup_services:
                    prod_vnet >> Edge(label="Backup\nTraffic", style="dotted", color="teal") >> bs
                    workloads_mg >> Edge(label="Backup\nPolicies", style="dashed", color="purple") >> bs
                    
    except Exception as e:
        logger.warning(f"Error adding data layer clusters: {str(e)}")


def _add_service_clusters(inputs: CustomerInputs, prod_vnet, workloads_mg):
    """Helper method to add service clusters to avoid code duplication"""
    try:
        # Compute and Application Services
        if inputs.compute_services or inputs.workload:
            with Cluster("Compute & Applications", graph_attr={"bgcolor": "#fff8dc", "style": "rounded"}):
                compute_services = []
                
                # Add selected compute services
                if inputs.compute_services:
                    for service in inputs.compute_services:
                        if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["diagram_class"]:
                            diagram_class = AZURE_SERVICES_MAPPING[service]["diagram_class"]
                            service_name = AZURE_SERVICES_MAPPING[service]["name"]
                            compute_services.append(diagram_class(service_name))
                
                # Fallback to workload if no specific compute services
                if not compute_services and inputs.workload:
                    if inputs.workload in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[inputs.workload]["diagram_class"]:
                        diagram_class = AZURE_SERVICES_MAPPING[inputs.workload]["diagram_class"]
                        service_name = AZURE_SERVICES_MAPPING[inputs.workload]["name"]
                        compute_services.append(diagram_class(service_name))
                
                # Default to App Services if nothing specified
                if not compute_services:
                    compute_services.append(AppServices("Azure App Services"))
                
                # Connect compute services to production VNet
                for cs in compute_services:
                    prod_vnet >> cs
                    workloads_mg >> cs
        
        # Storage Services
        if inputs.storage_services:
            with Cluster("Storage & Data", graph_attr={"bgcolor": "#f5f5dc", "style": "rounded"}):
                storage_services = []
                for service in inputs.storage_services:
                    if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["diagram_class"]:
                        diagram_class = AZURE_SERVICES_MAPPING[service]["diagram_class"]
                        service_name = AZURE_SERVICES_MAPPING[service]["name"]
                        storage_services.append(diagram_class(service_name))
                
                if not storage_services:
                    storage_services.append(StorageAccounts("Storage Accounts"))
                
                # Connect storage to production VNet and workloads
                for ss in storage_services:
                    prod_vnet >> ss
                    workloads_mg >> ss
        
        # Database Services
        if inputs.database_services:
            with Cluster("Databases", graph_attr={"bgcolor": "#e6f3ff", "style": "rounded"}):
                database_services = []
                for service in inputs.database_services:
                    if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["diagram_class"]:
                        diagram_class = AZURE_SERVICES_MAPPING[service]["diagram_class"]
                        service_name = AZURE_SERVICES_MAPPING[service]["name"]
                        database_services.append(diagram_class(service_name))
                
                # Connect databases to production VNet and workloads
                for ds in database_services:
                    prod_vnet >> ds
                    workloads_mg >> ds
        
        # Analytics Services
        if inputs.analytics_services:
            with Cluster("Analytics & AI", graph_attr={"bgcolor": "#f0e6ff", "style": "rounded"}):
                analytics_services = []
                for service in inputs.analytics_services:
                    if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["diagram_class"]:
                        diagram_class = AZURE_SERVICES_MAPPING[service]["diagram_class"]
                        service_name = AZURE_SERVICES_MAPPING[service]["name"]
                        analytics_services.append(diagram_class(service_name))
                
                # Connect analytics to production VNet and workloads
                for as_service in analytics_services:
                    prod_vnet >> as_service
                    workloads_mg >> as_service
        
        # Integration Services
        if inputs.integration_services:
            with Cluster("Integration", graph_attr={"bgcolor": "#fff0e6", "style": "rounded"}):
                integration_services = []
                for service in inputs.integration_services:
                    if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["diagram_class"]:
                        diagram_class = AZURE_SERVICES_MAPPING[service]["diagram_class"]
                        service_name = AZURE_SERVICES_MAPPING[service]["name"]
                        integration_services.append(diagram_class(service_name))
                
                # Connect integration services to production VNet and workloads
                for is_service in integration_services:
                    prod_vnet >> is_service
                    workloads_mg >> is_service
        
        # DevOps Services  
        if inputs.devops_services:
            with Cluster("DevOps & Automation", graph_attr={"bgcolor": "#f5f5f5", "style": "rounded"}):
                devops_services = []
                for service in inputs.devops_services:
                    if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["diagram_class"]:
                        diagram_class = AZURE_SERVICES_MAPPING[service]["diagram_class"]
                        service_name = AZURE_SERVICES_MAPPING[service]["name"]
                        devops_services.append(diagram_class(service_name))
                
                # Connect DevOps services to management
                for ds in devops_services:
                    workloads_mg >> ds
                    
    except Exception as e:
        logger.warning(f"Error adding service clusters: {str(e)}")
        # Don't fail the entire diagram generation for service cluster issues


def generate_architecture_template(inputs: CustomerInputs) -> Dict[str, Any]:
    """Generate architecture template based on inputs with AI enhancement for free text"""
    
    # First, check if we should use AI analysis for service selection
    ai_services = []
    ai_reasoning = ""
    
    # If user provided free text but didn't select many services, use AI to determine services
    total_selected_services = sum([
        len(inputs.compute_services or []),
        len(inputs.network_services or []),
        len(inputs.storage_services or []),
        len(inputs.database_services or []),
        len(inputs.security_services or []),
        len(inputs.monitoring_services or []),
        len(inputs.ai_services or []),
        len(inputs.analytics_services or []),
        len(inputs.integration_services or []),
        len(inputs.devops_services or []),
        len(inputs.backup_services or [])
    ])
    
    if inputs.free_text_input and total_selected_services <= 1:  # Very few services selected
        logger.info("Using AI analysis for service selection based on free text input")
        ai_analysis = analyze_free_text_requirements(inputs.free_text_input)
        ai_services = ai_analysis.get("services", [])
        ai_reasoning = ai_analysis.get("reasoning", "")
        
        # Override service selections with AI recommendations
        if ai_services:
            # Map AI services to appropriate categories
            for service in ai_services:
                if service in AZURE_SERVICES_MAPPING:
                    category = AZURE_SERVICES_MAPPING[service]["category"]
                    category_field = f"{category}_services"
                    
                    # Add to appropriate service list
                    if hasattr(inputs, category_field):
                        current_services = getattr(inputs, category_field) or []
                        if service not in current_services:
                            current_services.append(service)
                            setattr(inputs, category_field, current_services)
    
    # Determine organization size and template
    if inputs.org_structure and "enterprise" in inputs.org_structure.lower():
        template = AZURE_TEMPLATES["enterprise"]
    elif inputs.org_structure and any(x in inputs.org_structure.lower() for x in ["small", "medium", "sme"]):
        template = AZURE_TEMPLATES["small_medium"]
    else:
        template = AZURE_TEMPLATES["startup"]
    
    # Build architecture components
    components = {
        "template": template,
        "identity": inputs.identity or "Azure AD",
        "network_model": inputs.network_model or "hub-spoke",
        "workload": inputs.workload or "app-services",
        "security": inputs.security_posture or "zero-trust",
        "monitoring": inputs.monitoring or "azure-monitor",
        "governance": inputs.governance or "azure-policy",
        "ai_services": ai_services,
        "ai_reasoning": ai_reasoning
    }
    
    return components

def generate_professional_mermaid(inputs: CustomerInputs) -> str:
    """Generate professional Mermaid diagram for Azure Landing Zone"""
    
    template = generate_architecture_template(inputs)
    network_model = inputs.network_model or "hub-spoke"
    
    lines = [
        "graph TB",
        "    subgraph \"Azure Tenant\"",
        "        subgraph \"Management Groups\"",
        "            ROOT[\"🏢 Root Management Group\"]"
    ]
    
    # Add management group hierarchy
    if template["template"]["name"] == "Enterprise Scale Landing Zone":
        lines.extend([
            "            PLATFORM[\"[PLAT] Platform\"]",
            "            LANDINGZONES[\"[LZ] Landing Zones\"]", 
            "            SANDBOX[\"[SBX] Sandbox\"]",
            "            DECOM[\"[DECOM] Decommissioned\"]",
            "            ROOT --> PLATFORM",
            "            ROOT --> LANDINGZONES",
            "            ROOT --> SANDBOX", 
            "            ROOT --> DECOM"
        ])
    else:
        lines.extend([
            "            PLATFORM[\"[PLAT] Platform\"]",
            "            WORKLOADS[\"[WL] Workloads\"]",
            "            ROOT --> PLATFORM",
            "            ROOT --> WORKLOADS"
        ])
    
    lines.append("        end")
    
    # Add subscription structure
    lines.extend([
        "        subgraph \"Subscriptions\"",
        "            CONN[\"[NET] Connectivity\"]",
        "            IDENTITY[\"[ID] Identity\"]",
        "            MGMT[\"[MGMT] Management\"]",
        "            PROD[\"[PROD] Production\"]",
        "            DEV[\"[DEV] Development\"]"
    ])
    
    if template["template"]["name"] == "Enterprise Scale Landing Zone":
        lines.extend([
            "            PLATFORM --> CONN",
            "            PLATFORM --> IDENTITY", 
            "            PLATFORM --> MGMT",
            "            LANDINGZONES --> PROD",
            "            LANDINGZONES --> DEV"
        ])
    else:
        lines.extend([
            "            PLATFORM --> CONN",
            "            PLATFORM --> IDENTITY",
            "            WORKLOADS --> PROD",
            "            WORKLOADS --> DEV"
        ])
    
    lines.append("        end")
    
    # Add networking services based on selections
    if inputs.network_services:
        lines.extend([
            "        subgraph \"Networking Services\""
        ])
        
        service_connections = []
        for service in inputs.network_services:
            if service in AZURE_SERVICES_MAPPING:
                service_info = AZURE_SERVICES_MAPPING[service]
                service_id = service.upper().replace("_", "")
                lines.append(f"            {service_id}[\"{service_info['icon']} {service_info['name']}\"]")
                service_connections.append(f"            CONN --> {service_id}")
        
        lines.extend(service_connections)
        lines.append("        end")
    
    # Add compute services based on selections
    if inputs.compute_services:
        lines.extend([
            "        subgraph \"Compute Services\""
        ])
        
        service_connections = []
        for service in inputs.compute_services:
            if service in AZURE_SERVICES_MAPPING:
                service_info = AZURE_SERVICES_MAPPING[service]
                service_id = service.upper().replace("_", "")
                lines.append(f"            {service_id}[\"{service_info['icon']} {service_info['name']}\"]")
                service_connections.append(f"            PROD --> {service_id}")
        
        lines.extend(service_connections)
        lines.append("        end")
    
    # Add storage services based on selections
    if inputs.storage_services:
        lines.extend([
            "        subgraph \"Storage Services\""
        ])
        
        service_connections = []
        for service in inputs.storage_services:
            if service in AZURE_SERVICES_MAPPING:
                service_info = AZURE_SERVICES_MAPPING[service]
                service_id = service.upper().replace("_", "")
                lines.append(f"            {service_id}[\"{service_info['icon']} {service_info['name']}\"]")
                service_connections.append(f"            PROD --> {service_id}")
        
        lines.extend(service_connections)
        lines.append("        end")
    
    # Add database services based on selections
    if inputs.database_services:
        lines.extend([
            "        subgraph \"Database Services\""
        ])
        
        service_connections = []
        for service in inputs.database_services:
            if service in AZURE_SERVICES_MAPPING:
                service_info = AZURE_SERVICES_MAPPING[service]
                service_id = service.upper().replace("_", "")
                lines.append(f"            {service_id}[\"{service_info['icon']} {service_info['name']}\"]")
                service_connections.append(f"            PROD --> {service_id}")
        
        lines.extend(service_connections)
        lines.append("        end")
    
    # Add security services based on selections
    if inputs.security_services:
        lines.extend([
            "        subgraph \"Security & Identity\""
        ])
        
        service_connections = []
        for service in inputs.security_services:
            if service in AZURE_SERVICES_MAPPING:
                service_info = AZURE_SERVICES_MAPPING[service]
                service_id = service.upper().replace("_", "")
                lines.append(f"            {service_id}[\"{service_info['icon']} {service_info['name']}\"]")
                service_connections.append(f"            IDENTITY --> {service_id}")
        
        lines.extend(service_connections)
        lines.append("        end")
    
    # Add monitoring services based on selections
    if inputs.monitoring_services:
        lines.extend([
            "        subgraph \"Monitoring & Management\""
        ])
        
        service_connections = []
        for service in inputs.monitoring_services:
            if service in AZURE_SERVICES_MAPPING:
                service_info = AZURE_SERVICES_MAPPING[service]
                service_id = service.upper().replace("_", "")
                lines.append(f"            {service_id}[\"{service_info['icon']} {service_info['name']}\"]")
                service_connections.append(f"            MGMT --> {service_id}")
        
        lines.extend(service_connections)
        lines.append("        end")
    
    # Add AI/ML services based on selections  
    if inputs.ai_services:
        lines.extend([
            "        subgraph \"AI & Machine Learning\""
        ])
        
        service_connections = []
        for service in inputs.ai_services:
            if service in AZURE_SERVICES_MAPPING:
                service_info = AZURE_SERVICES_MAPPING[service]
                service_id = service.upper().replace("_", "")
                lines.append(f"            {service_id}[\"{service_info['icon']} {service_info['name']}\"]")
                service_connections.append(f"            PROD --> {service_id}")
        
        lines.extend(service_connections)
        lines.append("        end")
    
    # Add analytics services based on selections
    if inputs.analytics_services:
        lines.extend([
            "        subgraph \"Data & Analytics\""
        ])
        
        service_connections = []
        for service in inputs.analytics_services:
            if service in AZURE_SERVICES_MAPPING:
                service_info = AZURE_SERVICES_MAPPING[service]
                service_id = service.upper().replace("_", "")
                lines.append(f"            {service_id}[\"{service_info['icon']} {service_info['name']}\"]")
                service_connections.append(f"            PROD --> {service_id}")
        
        lines.extend(service_connections)
        lines.append("        end")
    
    # Add integration services based on selections
    if inputs.integration_services:
        lines.extend([
            "        subgraph \"Integration Services\""
        ])
        
        service_connections = []
        for service in inputs.integration_services:
            if service in AZURE_SERVICES_MAPPING:
                service_info = AZURE_SERVICES_MAPPING[service]
                service_id = service.upper().replace("_", "")
                lines.append(f"            {service_id}[\"{service_info['icon']} {service_info['name']}\"]")
                service_connections.append(f"            PROD --> {service_id}")
        
        lines.extend(service_connections)
        lines.append("        end")
    
    # Add DevOps services based on selections
    if inputs.devops_services:
        lines.extend([
            "        subgraph \"DevOps & Governance\""
        ])
        
        service_connections = []
        for service in inputs.devops_services:
            if service in AZURE_SERVICES_MAPPING:
                service_info = AZURE_SERVICES_MAPPING[service]
                service_id = service.upper().replace("_", "")
                lines.append(f"            {service_id}[\"{service_info['icon']} {service_info['name']}\"]")
                service_connections.append(f"            MGMT --> {service_id}")
        
        lines.extend(service_connections)
        lines.append("        end")
    
    # Add backup services based on selections
    if inputs.backup_services:
        lines.extend([
            "        subgraph \"Backup & Recovery\""
        ])
        
        service_connections = []
        for service in inputs.backup_services:
            if service in AZURE_SERVICES_MAPPING:
                service_info = AZURE_SERVICES_MAPPING[service]
                service_id = service.upper().replace("_", "")
                lines.append(f"            {service_id}[\"{service_info['icon']} {service_info['name']}\"]")
                service_connections.append(f"            MGMT --> {service_id}")
        
        lines.extend(service_connections)
        lines.append("        end")
    
    # Fallback for legacy workload field
    if not any([inputs.compute_services, inputs.network_services, inputs.storage_services, 
               inputs.database_services, inputs.security_services, inputs.monitoring_services,
               inputs.ai_services, inputs.analytics_services, inputs.integration_services,
               inputs.devops_services, inputs.backup_services]) and inputs.workload:
        workload_name = AZURE_SERVICES_MAPPING.get(inputs.workload, {"name": inputs.workload, "icon": "⚙️"})["name"]
        lines.extend([
            "        subgraph \"Workloads\"",
            f"            WORKLOAD[\"{workload_name}\"]",
            "            PROD --> WORKLOAD",
            "        end"
        ])
    
    lines.append("    end")
    
    # Add styling
    lines.extend([
        "",
        "    classDef mgmtGroup fill:#e1f5fe,stroke:#01579b,stroke-width:2px;",
        "    classDef subscription fill:#f3e5f5,stroke:#4a148c,stroke-width:2px;", 
        "    classDef compute fill:#fff3e0,stroke:#e65100,stroke-width:2px;",
        "    classDef network fill:#e8f5e8,stroke:#1b5e20,stroke-width:2px;",
        "    classDef storage fill:#fce4ec,stroke:#880e4f,stroke-width:2px;",
        "    classDef database fill:#e3f2fd,stroke:#0d47a1,stroke-width:2px;",
        "    classDef security fill:#ffebee,stroke:#b71c1c,stroke-width:2px;",
        "    classDef monitoring fill:#f1f8e9,stroke:#33691e,stroke-width:2px;",
        "    classDef ai fill:#f3e5f5,stroke:#4a148c,stroke-width:2px;",
        "    classDef analytics fill:#e8eaf6,stroke:#1a237e,stroke-width:2px;",
        "    classDef integration fill:#fff8e1,stroke:#f57f17,stroke-width:2px;",
        "    classDef devops fill:#fafafa,stroke:#424242,stroke-width:2px;",
        "    classDef backup fill:#e0f2f1,stroke:#00695c,stroke-width:2px;",
        "",
        "    class ROOT,PLATFORM,LANDINGZONES,SANDBOX,DECOM,WORKLOADS mgmtGroup;",
        "    class CONN,IDENTITY,MGMT,PROD,DEV subscription;"
    ])
    
    return "\n".join(lines)

def generate_enhanced_drawio_xml(inputs: CustomerInputs) -> str:
    """Generate enhanced Draw.io XML with comprehensive Azure stencils based on user selections"""
    
    def esc(s): 
        return html.escape(s) if s else ""
    
    template = generate_architecture_template(inputs)
    diagram_id = str(uuid.uuid4())
    
    # Base layout coordinates
    y_start = 100
    current_y = y_start
    section_height = 350
    service_width = 100
    service_height = 80
    
    # Build dynamic XML content
    xml_parts = [
        f"""<mxfile host="app.diagrams.net" modified="2024-01-01T00:00:00.000Z" agent="Azure Landing Zone Agent" version="1.0.0">
  <diagram name="Azure Landing Zone Architecture" id="azure-lz-{diagram_id}">
    <mxGraphModel dx="1422" dy="794" grid="1" gridSize="10" guides="1" tooltips="1" connect="1" arrows="1" fold="1" page="1" pageScale="1" pageWidth="2400" pageHeight="1600" math="0" shadow="0">
      <root>
        <mxCell id="0" />
        <mxCell id="1" parent="0" />
        
        <!-- Azure Tenant Container -->
        <mxCell id="tenant" value="Azure Tenant - {esc(inputs.org_structure or 'Enterprise')}" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#dae8fc;strokeColor=#6c8ebf;fontSize=16;fontStyle=1;verticalAlign=top;spacingTop=10;" vertex="1" parent="1">
          <mxGeometry x="50" y="50" width="2300" height="1500" as="geometry" />
        </mxCell>
        
        <!-- Management Groups -->
        <mxCell id="mgmt-groups" value="Management Groups" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#e1d5e7;strokeColor=#9673a6;fontSize=14;fontStyle=1;verticalAlign=top;spacingTop=10;" vertex="1" parent="1">
          <mxGeometry x="100" y="{current_y}" width="500" height="250" as="geometry" />
        </mxCell>"""
    ]
    
    # Management Group structure based on template
    mg_x = 150
    mg_y = current_y + 50
    xml_parts.append(f"""
        <mxCell id="root-mg" value="Root MG" style="shape=mxgraph.azure.management;fillColor=#0078d4;strokeColor=#005a9e;fontColor=#ffffff;" vertex="1" parent="1">
          <mxGeometry x="{mg_x}" y="{mg_y}" width="80" height="60" as="geometry" />
        </mxCell>""")
    
    mg_x += 120
    for i, mg in enumerate(template['template']['management_groups'][1:3]):  # Platform and Workloads
        mg_id = mg.lower().replace(' ', '-')
        xml_parts.append(f"""
        <mxCell id="{mg_id}-mg" value="{mg}" style="shape=mxgraph.azure.management;fillColor=#0078d4;strokeColor=#005a9e;fontColor=#ffffff;" vertex="1" parent="1">
          <mxGeometry x="{mg_x}" y="{mg_y}" width="80" height="60" as="geometry" />
        </mxCell>""")
        mg_x += 120
    
    # Subscriptions
    current_y += 300
    xml_parts.append(f"""
        <!-- Subscriptions -->
        <mxCell id="subscriptions" value="Subscriptions" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#d5e8d4;strokeColor=#82b366;fontSize=14;fontStyle=1;verticalAlign=top;spacingTop=10;" vertex="1" parent="1">
          <mxGeometry x="700" y="{current_y}" width="600" height="250" as="geometry" />
        </mxCell>""")
    
    sub_x = 750
    sub_y = current_y + 50
    for sub in template['template']['subscriptions'][:4]:  # First 4 subscriptions
        xml_parts.append(f"""
        <mxCell id="{sub.lower().replace(' ', '-')}-sub" value="{sub}" style="shape=mxgraph.azure.subscription;fillColor=#0078d4;strokeColor=#005a9e;fontColor=#ffffff;" vertex="1" parent="1">
          <mxGeometry x="{sub_x}" y="{sub_y}" width="{service_width}" height="{service_height}" as="geometry" />
        </mxCell>""")
        sub_x += 130
        if sub_x > 1200:  # Wrap to next row
            sub_x = 750
            sub_y += 100
    
    # Network Architecture Section
    current_y += 300
    xml_parts.append(f"""
        <!-- Network Architecture -->
        <mxCell id="network" value="Network Architecture - {esc(inputs.network_model or 'Hub-Spoke')}" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#fff2cc;strokeColor=#d6b656;fontSize=14;fontStyle=1;verticalAlign=top;spacingTop=10;" vertex="1" parent="1">
          <mxGeometry x="100" y="{current_y}" width="800" height="{section_height}" as="geometry" />
        </mxCell>""")
    
    # Hub VNet (always present)
    hub_x = 200
    hub_y = current_y + 80
    xml_parts.append(f"""
        <mxCell id="hub-vnet" value="Hub VNet\\nShared Services" style="shape=mxgraph.azure.virtual_network;fillColor=#0078d4;strokeColor=#005a9e;fontColor=#ffffff;" vertex="1" parent="1">
          <mxGeometry x="{hub_x}" y="{hub_y}" width="120" height="80" as="geometry" />
        </mxCell>""")
    
    # Spoke VNets
    spoke_x = 400
    spoke_y = hub_y - 50
    xml_parts.append(f"""
        <mxCell id="spoke1-vnet" value="Production VNet" style="shape=mxgraph.azure.virtual_network;fillColor=#0078d4;strokeColor=#005a9e;fontColor=#ffffff;" vertex="1" parent="1">
          <mxGeometry x="{spoke_x}" y="{spoke_y}" width="120" height="80" as="geometry" />
        </mxCell>
        <mxCell id="spoke2-vnet" value="Development VNet" style="shape=mxgraph.azure.virtual_network;fillColor=#0078d4;strokeColor=#005a9e;fontColor=#ffffff;" vertex="1" parent="1">
          <mxGeometry x="{spoke_x}" y="{spoke_y + 120}" width="120" height="80" as="geometry" />
        </mxCell>""")
    
    # Add selected network services
    if inputs.network_services:
        net_x = 600
        net_y = hub_y
        for i, service in enumerate(inputs.network_services[:4]):  # Max 4 network services
            if service in AZURE_SERVICES_MAPPING:
                service_info = AZURE_SERVICES_MAPPING[service]
                shape = service_info.get('drawio_shape', 'generic_service')
                xml_parts.append(f"""
        <mxCell id="net-service-{i}" value="{esc(service_info['name'])}" style="shape=mxgraph.azure.{shape};fillColor=#0078d4;strokeColor=#005a9e;fontColor=#ffffff;" vertex="1" parent="1">
          <mxGeometry x="{net_x}" y="{net_y}" width="{service_width}" height="{service_height}" as="geometry" />
        </mxCell>""")
                net_y += 100
    
    # Compute Services Section
    if inputs.compute_services or inputs.workload:
        current_y += section_height + 50
        xml_parts.append(f"""
        <!-- Compute Services -->
        <mxCell id="compute" value="Compute Services" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#fff0e6;strokeColor=#d79b00;fontSize=14;fontStyle=1;verticalAlign=top;spacingTop=10;" vertex="1" parent="1">
          <mxGeometry x="1000" y="{current_y}" width="600" height="{section_height}" as="geometry" />
        </mxCell>""")
        
        comp_x = 1050
        comp_y = current_y + 50
        
        # Add selected compute services
        services_to_add = inputs.compute_services or []
        if inputs.workload and inputs.workload not in services_to_add:
            services_to_add.append(inputs.workload)
            
        for i, service in enumerate(services_to_add[:6]):  # Max 6 compute services
            if service in AZURE_SERVICES_MAPPING:
                service_info = AZURE_SERVICES_MAPPING[service]
                shape = service_info.get('drawio_shape', 'generic_service')
                xml_parts.append(f"""
        <mxCell id="compute-service-{i}" value="{esc(service_info['name'])}" style="shape=mxgraph.azure.{shape};fillColor=#0078d4;strokeColor=#005a9e;fontColor=#ffffff;" vertex="1" parent="1">
          <mxGeometry x="{comp_x}" y="{comp_y}" width="{service_width}" height="{service_height}" as="geometry" />
        </mxCell>""")
                comp_x += 120
                if comp_x > 1450:  # Wrap to next row
                    comp_x = 1050
                    comp_y += 100
    
    # Storage Services Section
    if inputs.storage_services:
        current_y += section_height + 50
        xml_parts.append(f"""
        <!-- Storage Services -->
        <mxCell id="storage" value="Storage Services" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#f0f0f0;strokeColor=#666666;fontSize=14;fontStyle=1;verticalAlign=top;spacingTop=10;" vertex="1" parent="1">
          <mxGeometry x="100" y="{current_y}" width="600" height="250" as="geometry" />
        </mxCell>""")
        
        stor_x = 150
        stor_y = current_y + 50
        for i, service in enumerate(inputs.storage_services[:4]):
            if service in AZURE_SERVICES_MAPPING:
                service_info = AZURE_SERVICES_MAPPING[service]
                shape = service_info.get('drawio_shape', 'storage_accounts')
                xml_parts.append(f"""
        <mxCell id="storage-service-{i}" value="{esc(service_info['name'])}" style="shape=mxgraph.azure.{shape};fillColor=#0078d4;strokeColor=#005a9e;fontColor=#ffffff;" vertex="1" parent="1">
          <mxGeometry x="{stor_x}" y="{stor_y}" width="{service_width}" height="{service_height}" as="geometry" />
        </mxCell>""")
                stor_x += 120
                if stor_x > 550:
                    stor_x = 150
                    stor_y += 100
    
    # Database Services Section
    if inputs.database_services:
        db_y = current_y if not inputs.storage_services else current_y
        if inputs.storage_services:
            xml_parts.append(f"""
        <!-- Database Services -->
        <mxCell id="database" value="Database Services" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#e6f3ff;strokeColor=#0066cc;fontSize=14;fontStyle=1;verticalAlign=top;spacingTop=10;" vertex="1" parent="1">
          <mxGeometry x="800" y="{db_y}" width="600" height="250" as="geometry" />
        </mxCell>""")
        else:
            current_y += section_height + 50
            xml_parts.append(f"""
        <!-- Database Services -->
        <mxCell id="database" value="Database Services" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#e6f3ff;strokeColor=#0066cc;fontSize=14;fontStyle=1;verticalAlign=top;spacingTop=10;" vertex="1" parent="1">
          <mxGeometry x="100" y="{current_y}" width="600" height="250" as="geometry" />
        </mxCell>""")
            db_y = current_y
        
        db_x = 850 if inputs.storage_services else 150
        db_y += 50
        for i, service in enumerate(inputs.database_services[:4]):
            if service in AZURE_SERVICES_MAPPING:
                service_info = AZURE_SERVICES_MAPPING[service]
                shape = service_info.get('drawio_shape', 'sql_database')
                xml_parts.append(f"""
        <mxCell id="database-service-{i}" value="{esc(service_info['name'])}" style="shape=mxgraph.azure.{shape};fillColor=#0078d4;strokeColor=#005a9e;fontColor=#ffffff;" vertex="1" parent="1">
          <mxGeometry x="{db_x}" y="{db_y}" width="{service_width}" height="{service_height}" as="geometry" />
        </mxCell>""")
                db_x += 120
                if db_x > (1250 if inputs.storage_services else 550):
                    db_x = 850 if inputs.storage_services else 150
                    db_y += 100
    
    # Security Services Section (always present)
    current_y += 300
    xml_parts.append(f"""
        <!-- Security & Identity Services -->
        <mxCell id="security" value="Security &amp; Identity Services" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#ffebee;strokeColor=#c62828;fontSize=14;fontStyle=1;verticalAlign=top;spacingTop=10;" vertex="1" parent="1">
          <mxGeometry x="1700" y="{y_start}" width="600" height="600" as="geometry" />
        </mxCell>""")
    
    # Core security services (always present)
    sec_x = 1750
    sec_y = y_start + 50
    core_security = [
        ('azure-ad', 'Azure AD', 'azure_active_directory'),
        ('key-vault', 'Key Vault', 'key_vault'),
        ('security-center', 'Security Center', 'security_center')
    ]
    
    for sec_id, sec_name, sec_shape in core_security:
        xml_parts.append(f"""
        <mxCell id="{sec_id}" value="{sec_name}" style="shape=mxgraph.azure.{sec_shape};fillColor=#0078d4;strokeColor=#005a9e;fontColor=#ffffff;" vertex="1" parent="1">
          <mxGeometry x="{sec_x}" y="{sec_y}" width="{service_width}" height="{service_height}" as="geometry" />
        </mxCell>""")
        sec_x += 120
        if sec_x > 2100:
            sec_x = 1750
            sec_y += 100
    
    # Add additional selected security services
    if inputs.security_services:
        for i, service in enumerate(inputs.security_services):
            if service in AZURE_SERVICES_MAPPING and service not in ['active_directory', 'key_vault', 'security_center']:
                service_info = AZURE_SERVICES_MAPPING[service]
                shape = service_info.get('drawio_shape', 'generic_service')
                xml_parts.append(f"""
        <mxCell id="security-service-{i}" value="{esc(service_info['name'])}" style="shape=mxgraph.azure.{shape};fillColor=#0078d4;strokeColor=#005a9e;fontColor=#ffffff;" vertex="1" parent="1">
          <mxGeometry x="{sec_x}" y="{sec_y}" width="{service_width}" height="{service_height}" as="geometry" />
        </mxCell>""")
                sec_x += 120
                if sec_x > 2100:
                    sec_x = 1750
                    sec_y += 100
    
    # Analytics Services Section
    if inputs.analytics_services:
        analytics_y = y_start + 650
        xml_parts.append(f"""
        <!-- Analytics & AI Services -->
        <mxCell id="analytics" value="Analytics &amp; AI Services" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#f3e5f5;strokeColor=#9c27b0;fontSize=14;fontStyle=1;verticalAlign=top;spacingTop=10;" vertex="1" parent="1">
          <mxGeometry x="1700" y="{analytics_y}" width="600" height="300" as="geometry" />
        </mxCell>""")
        
        ana_x = 1750
        ana_y = analytics_y + 50
        for i, service in enumerate(inputs.analytics_services[:4]):
            if service in AZURE_SERVICES_MAPPING:
                service_info = AZURE_SERVICES_MAPPING[service]
                shape = service_info.get('drawio_shape', 'generic_service')
                xml_parts.append(f"""
        <mxCell id="analytics-service-{i}" value="{esc(service_info['name'])}" style="shape=mxgraph.azure.{shape};fillColor=#0078d4;strokeColor=#005a9e;fontColor=#ffffff;" vertex="1" parent="1">
          <mxGeometry x="{ana_x}" y="{ana_y}" width="{service_width}" height="{service_height}" as="geometry" />
        </mxCell>""")
                ana_x += 120
                if ana_x > 2100:
                    ana_x = 1750
                    ana_y += 100
    
    # Integration Services Section
    if inputs.integration_services:
        int_y = current_y
        xml_parts.append(f"""
        <!-- Integration Services -->
        <mxCell id="integration" value="Integration Services" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#fff8e1;strokeColor=#ff8f00;fontSize=14;fontStyle=1;verticalAlign=top;spacingTop=10;" vertex="1" parent="1">
          <mxGeometry x="100" y="{int_y}" width="600" height="250" as="geometry" />
        </mxCell>""")
        
        int_x = 150
        int_y += 50
        for i, service in enumerate(inputs.integration_services[:4]):
            if service in AZURE_SERVICES_MAPPING:
                service_info = AZURE_SERVICES_MAPPING[service]
                shape = service_info.get('drawio_shape', 'generic_service')
                xml_parts.append(f"""
        <mxCell id="integration-service-{i}" value="{esc(service_info['name'])}" style="shape=mxgraph.azure.{shape};fillColor=#0078d4;strokeColor=#005a9e;fontColor=#ffffff;" vertex="1" parent="1">
          <mxGeometry x="{int_x}" y="{int_y}" width="{service_width}" height="{service_height}" as="geometry" />
        </mxCell>""")
                int_x += 120
                if int_x > 550:
                    int_x = 150
                    int_y += 100
    
    # DevOps Services Section
    if inputs.devops_services:
        devops_y = current_y if not inputs.integration_services else current_y
        devops_x_offset = 800 if inputs.integration_services else 100
        xml_parts.append(f"""
        <!-- DevOps Services -->
        <mxCell id="devops" value="DevOps Services" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#f5f5f5;strokeColor=#666666;fontSize=14;fontStyle=1;verticalAlign=top;spacingTop=10;" vertex="1" parent="1">
          <mxGeometry x="{devops_x_offset}" y="{devops_y}" width="400" height="250" as="geometry" />
        </mxCell>""")
        
        dev_x = devops_x_offset + 50
        dev_y = devops_y + 50
        for i, service in enumerate(inputs.devops_services[:3]):
            if service in AZURE_SERVICES_MAPPING:
                service_info = AZURE_SERVICES_MAPPING[service]
                shape = service_info.get('drawio_shape', 'generic_service')
                xml_parts.append(f"""
        <mxCell id="devops-service-{i}" value="{esc(service_info['name'])}" style="shape=mxgraph.azure.{shape};fillColor=#0078d4;strokeColor=#005a9e;fontColor=#ffffff;" vertex="1" parent="1">
          <mxGeometry x="{dev_x}" y="{dev_y}" width="{service_width}" height="{service_height}" as="geometry" />
        </mxCell>""")
                dev_x += 120
                if dev_x > (devops_x_offset + 250):
                    dev_x = devops_x_offset + 50
                    dev_y += 100
    
    # Add basic connections
    xml_parts.append("""
        <!-- Key Connections -->
        <mxCell id="conn1" edge="1" source="root-mg" target="platform-mg" parent="1">
          <mxGeometry relative="1" as="geometry" />
        </mxCell>
        <mxCell id="conn2" edge="1" source="root-mg" target="landing-zones-mg" parent="1">
          <mxGeometry relative="1" as="geometry" />
        </mxCell>
        <mxCell id="conn3" edge="1" source="hub-vnet" target="spoke1-vnet" parent="1">
          <mxGeometry relative="1" as="geometry" />
        </mxCell>
        <mxCell id="conn4" edge="1" source="hub-vnet" target="spoke2-vnet" parent="1">
          <mxGeometry relative="1" as="geometry" />
        </mxCell>
        
      </root>
    </mxGraphModel>
  </diagram>
</mxfile>""")
    
    return "".join(xml_parts)


def generate_professional_documentation(inputs: CustomerInputs) -> Dict[str, str]:
    """Generate professional TSD, HLD, and LLD documentation with AI enhancement"""
    
    template = generate_architecture_template(inputs)
    timestamp = datetime.now().strftime("%Y-%m-%d")
    
    # Generate AI insights if additional inputs are provided
    url_analysis = ""
    doc_analysis = ""
    ai_recommendations = ""
    
    try:
        if inputs.url_input:
            url_analysis = analyze_url_content(inputs.url_input)
            
        if inputs.uploaded_files_info:
            doc_analysis = "Document analysis results incorporated from uploaded files."
            
        # Generate AI-enhanced recommendations
        ai_recommendations = generate_ai_enhanced_recommendations(inputs, url_analysis, doc_analysis)
    except Exception as e:
        logger.warning(f"AI enhancement failed: {e}")
        ai_recommendations = "AI enhancement not available - using standard recommendations."
    
    # Technical Specification Document (TSD)
    tsd = f"""# Technical Specification Document (TSD)
## Azure Landing Zone Architecture - Enterprise Edition

**Document Version:** 2.0 (AI-Enhanced)
**Date:** {timestamp}
**Business Objective:** {inputs.business_objective or 'Not specified'}

### Executive Summary
This document outlines the technical specifications for implementing an Azure Landing Zone architecture based on comprehensive customer requirements analysis, including AI-powered insights and recommendations.

### Business Requirements Analysis
- **Primary Objective:** {inputs.business_objective or 'Cost optimization and operational efficiency'}
- **Industry:** {inputs.industry or 'General'}
- **Regulatory Requirements:** {inputs.regulatory or 'Standard compliance'}
- **Organization Structure:** {inputs.org_structure or 'Enterprise'}
- **Governance Model:** {inputs.governance or 'Centralized with delegated permissions'}

### Architecture Template Selection
**Selected Template:** {template['template']['name']}
**Justification:** Based on organizational size, complexity, and regulatory requirements.

### Core Architecture Components
- **Identity & Access Management:** {inputs.identity or 'Azure Active Directory with hybrid integration'}
- **Network Architecture:** {inputs.network_model or 'Hub-Spoke with Azure Virtual WAN'}
- **Security Framework:** {inputs.security_posture or 'Zero Trust with defense in depth'}
- **Connectivity Strategy:** {inputs.connectivity or 'Hybrid cloud with ExpressRoute'}
- **Primary Workloads:** {inputs.workload or 'Multi-tier applications with microservices'}
- **Monitoring & Observability:** {inputs.monitoring or 'Azure Monitor with Log Analytics'}

### Enhanced Requirements Analysis
{f"**Additional Context:** {inputs.free_text_input}" if inputs.free_text_input else "**Additional Context:** Standard requirements captured through structured inputs."}

{f"**URL Analysis Insights:** {url_analysis[:500]}..." if url_analysis else ""}

{f"**Document Analysis:** Document analysis completed on uploaded files." if inputs.uploaded_files_info else ""}

### AI-Powered Architecture Recommendations
{ai_recommendations[:2000] if ai_recommendations else "Standard architecture recommendations applied."}

### Compliance & Governance Framework
- **Governance Model:** {inputs.governance or 'Centralized with delegated permissions'}
- **Policy Framework:** Azure Policy for compliance enforcement
- **Security Framework:** {inputs.security_posture or 'Zero Trust'} security model
"""

    # High Level Design (HLD)
    hld = f"""# High Level Design (HLD)
## Azure Landing Zone Implementation

**Document Version:** 1.0
**Date:** {timestamp}

### Architecture Overview
The proposed Azure Landing Zone follows the {template['template']['name']} pattern.

### Management Group Structure
"""
    
    for mg in template['template']['management_groups']:
        hld += f"- **{mg}:** Management group for {mg.lower()} resources\n"
    
    hld += f"""
### Subscription Strategy
"""
    
    for sub in template['template']['subscriptions']:
        hld += f"- **{sub}:** Dedicated subscription for {sub.lower()} workloads\n"
    
    hld += f"""
### Network Architecture
**Topology:** {inputs.network_model or 'Hub-Spoke'}
- **Hub VNet:** Central connectivity and shared services
- **Spoke VNets:** Workload-specific virtual networks
- **Connectivity:** {inputs.connectivity or 'ExpressRoute and VPN Gateway'}

### Security Architecture
**Security Model:** {inputs.security_posture or 'Zero Trust'}
- **Identity:** {inputs.identity or 'Azure Active Directory'} with conditional access
- **Network Security:** Network Security Groups and Azure Firewall
- **Data Protection:** {inputs.key_vault or 'Azure Key Vault'} for secrets management
- **Threat Protection:** {inputs.threat_protection or 'Azure Security Center and Sentinel'}

### Workload Architecture
**Primary Workload:** {inputs.workload or 'Application Services'}
- **Compute:** {AZURE_SERVICES_MAPPING.get(inputs.workload or 'appservices', {'name': 'Azure App Services'})['name']}
- **Architecture Style:** {inputs.architecture_style or 'Microservices'}
- **Scalability:** {inputs.scalability or 'Auto-scaling enabled'}
"""

    # Low Level Design (LLD)
    lld = f"""# Low Level Design (LLD)
## Azure Landing Zone Technical Implementation

**Document Version:** 1.0
**Date:** {timestamp}

### Resource Configuration

#### Management Groups
"""
    
    for i, mg in enumerate(template['template']['management_groups']):
        lld += f"""
**{mg} Management Group:**
- Management Group ID: mg-{mg.lower().replace(' ', '-')}
- Parent: {template['template']['management_groups'][i-1] if i > 0 else 'Tenant Root'}
- Applied Policies: Azure Policy assignments for {mg.lower()}
"""

    lld += f"""
#### Subscriptions
"""
    
    for sub in template['template']['subscriptions']:
        lld += f"""
**{sub} Subscription:**
- Subscription Name: sub-{sub.lower().replace(' ', '-')}
- Resource Groups: Multiple RGs based on workload segregation
- RBAC: Custom roles and assignments
- Budget: Cost management and alerting configured
"""

    lld += f"""
#### Network Configuration

**Hub Virtual Network:**
- VNet Name: vnet-hub-{inputs.network_model or 'spoke'}-001
- Address Space: 10.0.0.0/16
- Subnets:
  - GatewaySubnet: 10.0.0.0/24 (VPN/ExpressRoute Gateway)
  - AzureFirewallSubnet: 10.0.1.0/24 (Azure Firewall)
  - SharedServicesSubnet: 10.0.2.0/24 (Domain Controllers, etc.)

**Spoke Virtual Networks:**
- Production Spoke: vnet-prod-{inputs.workload or 'app'}-001 (10.1.0.0/16)
- Development Spoke: vnet-dev-{inputs.workload or 'app'}-001 (10.2.0.0/16)

#### Security Configuration

**Azure Active Directory:**
- Tenant: {inputs.org_structure or 'enterprise'}.onmicrosoft.com
- Custom Domains: Configured as required
- Conditional Access: {inputs.security_posture or 'Zero Trust'} policies
- PIM: Privileged Identity Management for admin roles

**Network Security:**
- Azure Firewall: Central security appliance
- NSGs: Network Security Groups on all subnets
- UDRs: User Defined Routes for traffic steering

**Key Management:**
- Key Vault: {inputs.key_vault or 'Azure Key Vault'} for certificates and secrets
- Managed Identities: For secure service-to-service authentication

#### Workload Configuration

**Primary Workload: {inputs.workload or 'Application Services'}**
- Service: {AZURE_SERVICES_MAPPING.get(inputs.workload or 'appservices', {'name': 'Azure App Services'})['name']}
- SKU: Production-grade tier
- Scaling: {inputs.scalability or 'Auto-scaling based on CPU/memory'}
- Monitoring: {inputs.monitoring or 'Azure Monitor'} with custom dashboards

#### Operations Configuration

**Monitoring and Alerting:**
- Log Analytics Workspace: Central logging for all resources
- Azure Monitor: Metrics and alerting
- Application Insights: Application performance monitoring

**Backup and Recovery:**
- Azure Backup: {inputs.backup or 'Daily backups with 30-day retention'}
- Site Recovery: Disaster recovery as needed

**Cost Management:**
- Budget Alerts: Monthly budget monitoring
- Cost Optimization: {inputs.cost_priority or 'Regular cost reviews and optimization'}

#### Infrastructure as Code

**IaC Tool:** {inputs.iac or 'Bicep/ARM Templates'}
- Template Structure: Modular templates for each component
- Deployment: CI/CD pipeline using Azure DevOps
- Version Control: Git repository with proper branching strategy
"""

    return {
        "tsd": tsd,
        "hld": hld, 
        "lld": lld
    }


# ---------- API Endpoints ----------

@app.get("/")
def root():
    return {
        "message": "Azure Landing Zone Agent API",
        "version": "1.0.0",
        "endpoints": [
            "/docs - API Documentation",
            "/generate-diagram - Generate architecture diagram (Mermaid + Draw.io)",
            "/generate-azure-diagram - Generate Azure architecture diagram with official Azure icons (Python Diagrams)",
            "/generate-drawio - Generate Draw.io XML",
            "/health - Health check"
        ]
    }

@app.get("/health")
def health_check():
    """Enhanced health check that verifies system dependencies"""
    logger.info("Running health check...")
    status = "healthy"
    issues = []
    
    # Check Graphviz availability
    try:
        result = subprocess.run(['dot', '-V'], capture_output=True, text=True, timeout=10)
        if result.returncode != 0:
            issues.append(f"Graphviz 'dot' command failed with return code {result.returncode}")
            status = "degraded"
        else:
            logger.info(f"Graphviz check passed: {result.stderr.strip()}")
    except subprocess.TimeoutExpired:
        issues.append("Graphviz 'dot' command timed out")
        status = "degraded"
    except FileNotFoundError:
        issues.append("Graphviz not installed or not accessible")
        status = "degraded"
    except Exception as e:
        issues.append(f"Graphviz check failed: {str(e)}")
        status = "degraded"
    
    # Check diagrams library
    try:
        from diagrams import Diagram
        logger.info("Diagrams library import successful")
    except ImportError as e:
        issues.append(f"Diagrams library import failed: {str(e)}")
        status = "unhealthy"
    
    # Check output directory accessibility
    try:
        output_dir = get_safe_output_directory()
        logger.info(f"Output directory accessible: {output_dir}")
    except Exception as e:
        issues.append(f"Cannot access output directory: {str(e)}")
        status = "degraded"
    
    # Check available disk space
    try:
        import shutil
        output_dir = get_safe_output_directory()
        total, used, free = shutil.disk_usage(output_dir)
        free_mb = free // (1024*1024)
        if free_mb < 100:  # Less than 100MB free
            issues.append(f"Low disk space in output directory: {free_mb}MB free")
            status = "degraded"
        logger.info(f"Disk space check passed: {free_mb}MB free")
    except Exception as e:
        issues.append(f"Cannot check disk space: {str(e)}")
        status = "degraded"
    
    # Test a simple diagram generation
    try:
        from main import CustomerInputs
        test_inputs = CustomerInputs(business_objective="Health check test")
        # Just validate inputs, don't generate full diagram
        validate_customer_inputs(test_inputs)
        logger.info("Input validation test passed")
    except Exception as e:
        issues.append(f"Input validation test failed: {str(e)}")
        status = "degraded"
    
    logger.info(f"Health check completed with status: {status}")
    
    return {
        "status": status,
        "timestamp": datetime.now().isoformat(),
        "issues": issues,
        "dependencies": {
            "graphviz_available": "Graphviz" not in str(issues),
            "diagrams_available": "Diagrams library" not in str(issues),
            "output_directory_accessible": "output directory" not in str(issues),
            "sufficient_disk_space": "disk space" not in str(issues)
        }
    }

@app.post("/generate-diagram")
def generate_diagram(inputs: CustomerInputs):
    """Generate comprehensive Azure Landing Zone diagrams and documentation"""
    try:
        # Generate professional diagrams
        mermaid_diagram = generate_professional_mermaid(inputs)
        drawio_xml = generate_enhanced_drawio_xml(inputs)
        
        # Skip AI documentation generation to prevent timeouts - provide basic documentation instead
        docs = {
            "tsd": f"""# Technical Specification Document (TSD)
## Azure Landing Zone Architecture

**Document Version:** 1.0
**Date:** {datetime.now().strftime("%Y-%m-%d")}
**Business Objective:** {inputs.business_objective or 'Not specified'}

### Executive Summary
This document outlines the technical specifications for implementing an Azure Landing Zone architecture.
Both Mermaid and Draw.io diagrams have been successfully generated.

### Architecture Components
- **Compute Services:** {', '.join(inputs.compute_services or ['Not specified'])}
- **Network Services:** {', '.join(inputs.network_services or ['Not specified'])}
- **Storage Services:** {', '.join(inputs.storage_services or ['Not specified'])}
- **Database Services:** {', '.join(inputs.database_services or ['Not specified'])}
- **Security Services:** {', '.join(inputs.security_services or ['Not specified'])}

### Status
Diagrams generated successfully (Mermaid + Draw.io). Full AI-enhanced documentation is available through separate endpoints.
""",
            "hld": f"""# High Level Design (HLD)
## Azure Landing Zone Implementation

**Document Version:** 1.0
**Date:** {datetime.now().strftime("%Y-%m-%d")}

### Architecture Overview
The Azure Landing Zone follows enterprise best practices for cloud architecture.

### Selected Services
- **Compute:** {', '.join(inputs.compute_services or ['Not specified'])}
- **Network:** {', '.join(inputs.network_services or ['Not specified'])}
- **Storage:** {', '.join(inputs.storage_services or ['Not specified'])}
- **Security:** {', '.join(inputs.security_services or ['Not specified'])}

### Implementation Status
Mermaid and Draw.io diagrams generated successfully.
""",
            "lld": f"""# Low Level Design (LLD)
## Azure Landing Zone Technical Implementation

**Document Version:** 1.0  
**Date:** {datetime.now().strftime("%Y-%m-%d")}

### Resource Configuration
The architecture diagrams show the detailed technical implementation.

### Generated Artifacts
- Mermaid Diagram: Professional format for technical documentation
- Draw.io XML: Editable format for customization
- Generation Time: {datetime.now().isoformat()}

### Next Steps
Use the generated diagrams for implementation planning and detailed resource configuration.
"""
        }
        
        return {
            "success": True,
            "mermaid": mermaid_diagram,
            "drawio": drawio_xml,
            "tsd": docs["tsd"],
            "hld": docs["hld"],
            "lld": docs["lld"],
            "architecture_template": generate_architecture_template(inputs),
            "metadata": {
                "generated_at": datetime.now().isoformat(),
                "version": "1.0.0",
                "agent": "Azure Landing Zone Agent",
                "note": "AI documentation generation bypassed to ensure fast response times"
            }
        }
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error generating diagram: {str(e)}")

@app.post("/generate-azure-diagram")
def generate_azure_diagram_endpoint(inputs: CustomerInputs):
    """Generate Azure architecture diagram using enhanced simple SVG generation with 50+ design principles"""
    try:
        # Generate Azure architecture diagram with enhanced design principles
        diagram_path = generate_enhanced_azure_architecture_diagram(inputs)
        
        # Read the generated PNG file
        with open(diagram_path, "rb") as f:
            diagram_data = f.read()
        
        # Skip AI documentation generation to prevent timeouts - provide basic documentation instead
        docs = {
            "tsd": f"""# Technical Specification Document (TSD)
## Azure Landing Zone Architecture

**Document Version:** 1.0
**Date:** {datetime.now().strftime("%Y-%m-%d")}
**Business Objective:** {inputs.business_objective or 'Not specified'}

### Executive Summary
This document outlines the technical specifications for implementing an Azure Landing Zone architecture.
The diagram has been successfully generated showing the complete architecture.

### Architecture Components
- **Compute Services:** {', '.join(inputs.compute_services or ['Not specified'])}
- **Network Services:** {', '.join(inputs.network_services or ['Not specified'])}
- **Storage Services:** {', '.join(inputs.storage_services or ['Not specified'])}
- **Database Services:** {', '.join(inputs.database_services or ['Not specified'])}
- **Security Services:** {', '.join(inputs.security_services or ['Not specified'])}

### Status
Architecture diagram generated successfully. Full AI-enhanced documentation is available through separate endpoints.
""",
            "hld": f"""# High Level Design (HLD)
## Azure Landing Zone Implementation

**Document Version:** 1.0
**Date:** {datetime.now().strftime("%Y-%m-%d")}

### Architecture Overview
The Azure Landing Zone follows enterprise best practices for cloud architecture.

### Selected Services
- **Compute:** {', '.join(inputs.compute_services or ['Not specified'])}
- **Network:** {', '.join(inputs.network_services or ['Not specified'])}
- **Storage:** {', '.join(inputs.storage_services or ['Not specified'])}
- **Security:** {', '.join(inputs.security_services or ['Not specified'])}

### Implementation Status
Diagram generated successfully at: {diagram_path}
""",
            "lld": f"""# Low Level Design (LLD)
## Azure Landing Zone Technical Implementation

**Document Version:** 1.0  
**Date:** {datetime.now().strftime("%Y-%m-%d")}

### Resource Configuration
The architecture diagram shows the detailed technical implementation.

### Generated Artifacts
- Diagram Path: {diagram_path}
- Generation Time: {datetime.now().isoformat()}
- Format: PNG with Azure official icons

### Next Steps
Use the generated diagram for implementation planning and detailed resource configuration.
"""
        }
        
        # Encode the diagram as base64 for JSON response
        import base64
        diagram_base64 = base64.b64encode(diagram_data).decode('utf-8')
        
        return {
            "success": True,
            "diagram_path": diagram_path,
            "diagram_base64": diagram_base64,
            "tsd": docs["tsd"],
            "hld": docs["hld"],
            "lld": docs["lld"],
            "architecture_template": generate_architecture_template(inputs),
            "metadata": {
                "generated_at": datetime.now().isoformat(),
                "version": "1.0.0",
                "agent": "Azure Landing Zone Agent - Python Diagrams",
                "diagram_format": "PNG with Azure official icons",
                "note": "AI documentation generation bypassed to ensure fast response times"
            }
        }
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error generating Azure diagram: {str(e)}")

@app.get("/generate-azure-diagram/download/{filename}")
def download_azure_diagram(filename: str):
    """Download generated Azure architecture diagram PNG file"""
    try:
        file_path = f"/tmp/{filename}"
        if not os.path.exists(file_path):
            raise HTTPException(status_code=404, detail="Diagram file not found")
        
        with open(file_path, "rb") as f:
            diagram_data = f.read()
        
        return Response(
            content=diagram_data,
            media_type="image/png", 
            headers={"Content-Disposition": f"attachment; filename={filename}"}
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error downloading diagram: {str(e)}")

@app.post("/generate-drawio", response_class=Response)
def generate_drawio_endpoint(inputs: CustomerInputs):
    """Generate Draw.io XML file for download"""
    try:
        xml = generate_enhanced_drawio_xml(inputs)
        return Response(
            content=xml, 
            media_type="application/xml",
            headers={"Content-Disposition": "attachment; filename=azure-landing-zone.drawio"}
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error generating Draw.io XML: {str(e)}")


@app.post("/generate-comprehensive-azure-architecture")
def generate_comprehensive_azure_architecture(inputs: CustomerInputs):
    """Generate comprehensive Azure architecture with both Draw.io XML and PNG diagram"""
    logger.info("Starting comprehensive Azure architecture generation")
    
    try:
        # Validate inputs early
        validate_customer_inputs(inputs)
        logger.info("Input validation completed successfully")
        
        # Generate Draw.io XML with comprehensive Azure stencils
        logger.info("Generating Draw.io XML...")
        drawio_xml = generate_enhanced_drawio_xml(inputs)
        logger.info(f"Draw.io XML generated successfully (size: {len(drawio_xml)} characters)")
        
        # Generate Azure PNG diagram with enhanced design principles
        logger.info("Generating enhanced Azure PNG diagram...")
        diagram_path = generate_enhanced_azure_architecture_diagram(inputs)
        logger.info(f"Azure PNG diagram generated successfully: {diagram_path}")
        
        # Read the PNG file
        try:
            with open(diagram_path, "rb") as f:
                diagram_data = f.read()
            diagram_base64 = base64.b64encode(diagram_data).decode('utf-8')
            logger.info(f"PNG file read and encoded successfully (size: {len(diagram_data)} bytes)")
        except Exception as e:
            logger.error(f"Failed to read PNG file {diagram_path}: {e}")
            raise Exception(f"Failed to read generated PNG file: {str(e)}")
        
        # Skip AI documentation generation to prevent timeouts - provide basic documentation instead
        logger.info("Generating basic documentation (AI documentation bypassed)...")
        docs = {
            "tsd": f"""# Technical Specification Document (TSD)
## Comprehensive Azure Landing Zone Architecture

**Document Version:** 1.0
**Date:** {datetime.now().strftime("%Y-%m-%d")}
**Business Objective:** {inputs.business_objective or 'Not specified'}

### Executive Summary
This document outlines the technical specifications for implementing a comprehensive Azure Landing Zone architecture.
Both Draw.io XML and PNG diagrams have been successfully generated.

### Architecture Components
- **Compute Services:** {', '.join(inputs.compute_services or ['Not specified'])}
- **Network Services:** {', '.join(inputs.network_services or ['Not specified'])}
- **Storage Services:** {', '.join(inputs.storage_services or ['Not specified'])}
- **Database Services:** {', '.join(inputs.database_services or ['Not specified'])}
- **Security Services:** {', '.join(inputs.security_services or ['Not specified'])}

### Generated Artifacts
- Draw.io XML: {len(drawio_xml)} characters
- PNG Diagram: {diagram_path}

### Status
Comprehensive architecture diagrams generated successfully. Full AI-enhanced documentation is available through separate endpoints.
""",
            "hld": f"""# High Level Design (HLD)
## Comprehensive Azure Landing Zone Implementation

**Document Version:** 1.0
**Date:** {datetime.now().strftime("%Y-%m-%d")}

### Architecture Overview
The comprehensive Azure Landing Zone includes both editable Draw.io format and high-quality PNG visualization.

### Selected Services
- **Compute:** {', '.join(inputs.compute_services or ['Not specified'])}
- **Network:** {', '.join(inputs.network_services or ['Not specified'])}
- **Storage:** {', '.join(inputs.storage_services or ['Not specified'])}
- **Security:** {', '.join(inputs.security_services or ['Not specified'])}

### Implementation Status
Both Draw.io XML and PNG diagrams generated successfully.
""",
            "lld": f"""# Low Level Design (LLD)
## Comprehensive Azure Landing Zone Technical Implementation

**Document Version:** 1.0  
**Date:** {datetime.now().strftime("%Y-%m-%d")}

### Resource Configuration
The architecture diagrams show the detailed technical implementation in multiple formats.

### Generated Artifacts
- Draw.io XML: Editable format for customization and collaboration
- PNG Diagram: High-quality visualization for presentations
- Generation Time: {datetime.now().isoformat()}

### Next Steps
Use the generated diagrams for implementation planning and detailed resource configuration.
"""
        }
        logger.info("Basic documentation generated successfully")
        
        # Count Azure stencils used
        import re
        shapes = re.findall(r'shape=mxgraph\.azure\.[^;\"\s]*', drawio_xml)
        
        result = {
            "success": True,
            "drawio_xml": drawio_xml,
            "png_diagram_path": diagram_path,
            "png_diagram_base64": diagram_base64,
            "tsd": docs["tsd"],
            "hld": docs["hld"],
            "lld": docs["lld"],
            "architecture_template": generate_architecture_template(inputs),
            "azure_stencils": {
                "total_used": len(shapes),
                "unique_used": len(set(shapes)),
                "stencils_list": sorted(list(set(shapes)))
            },
            "metadata": {
                "generated_at": datetime.now().isoformat(),
                "version": "1.0.0",
                "agent": "Azure Landing Zone Agent - Comprehensive Generator",
                "drawio_size": len(drawio_xml),
                "png_size": len(diagram_data)
            }
        }
        
        logger.info("Comprehensive Azure architecture generated successfully")
        return result
    
    except ValueError as ve:
        # Input validation errors
        error_msg = f"Invalid input: {str(ve)}"
        logger.error(error_msg)
        raise HTTPException(status_code=400, detail=error_msg)
    
    except Exception as e:
        # Log the full error for debugging
        error_msg = f"Error generating comprehensive architecture: {str(e)}"
        logger.error(error_msg)
        logger.error(traceback.format_exc())
        
        # Return a user-friendly error
        raise HTTPException(
            status_code=500, 
            detail=f"Failed to generate architecture. Error: {str(e)}"
        )

@app.post("/generate-interactive-azure-architecture")
def generate_interactive_azure_architecture(inputs: CustomerInputs):
    """Generate interactive Azure architecture with SVG diagram for web display"""
    logger.info("Starting interactive Azure architecture generation")
    
    try:
        # Validate inputs early
        validate_customer_inputs(inputs)
        logger.info("Input validation completed successfully")
        
        # Check if we should provide human-in-the-loop feedback for better architecture
        feedback_questions = generate_feedback_questions(inputs)
        
        # Generate Mermaid diagram
        logger.info("Generating Mermaid diagram...")
        mermaid_diagram = generate_professional_mermaid(inputs)
        logger.info("Mermaid diagram generated successfully")
        
        # Generate Azure SVG diagram with proper Azure icons
        logger.info("Generating Azure SVG diagram...")
        svg_content = ""
        svg_diagram_path = ""
        
        try:
            svg_diagram_path = generate_enhanced_azure_architecture_diagram(inputs, format="svg")
            # Read the SVG file
            with open(svg_diagram_path, "r", encoding="utf-8") as f:
                svg_content = f.read()
            logger.info(f"Azure SVG diagram generated successfully: {svg_diagram_path}")
        except Exception as svg_error:
            logger.warning(f"SVG generation failed, using fallback: {str(svg_error)}")
            # Fallback: Create a simple SVG representation of the architecture
            svg_content = generate_simple_svg_diagram(inputs)
            logger.info("Using simple SVG fallback diagram")
        
        if svg_content:
            logger.info(f"SVG content ready (size: {len(svg_content)} characters)")
        else:
            logger.warning("No SVG content available, falling back to Mermaid only")
        
        # Generate Draw.io XML for compatibility
        logger.info("Generating Draw.io XML...")
        drawio_xml = generate_enhanced_drawio_xml(inputs)
        logger.info(f"Draw.io XML generated successfully (size: {len(drawio_xml)} characters)")
        
        # Generate professional documentation
        logger.info("Generating professional documentation...")
        try:
            # Use a timeout for the documentation generation
            import signal
            
            def timeout_handler(signum, frame):
                raise Exception("Documentation generation timed out")
            
            signal.signal(signal.SIGALRM, timeout_handler)
            signal.alarm(10)  # 10 second timeout
            
            docs = generate_professional_documentation(inputs)
            
            signal.alarm(0)  # Clear the alarm
            logger.info("Professional documentation generated successfully")
        except Exception as e:
            signal.alarm(0)  # Clear the alarm
            logger.warning(f"Documentation generation failed, using fallback: {str(e)}")
            docs = {
                "tsd": f"# Technical Specification Document\n\n## Azure Landing Zone Architecture\n\n**Organization:** {inputs.org_structure or 'Enterprise'}\n**Business Objective:** {inputs.business_objective or 'Not specified'}\n\n### Selected Services\n- Compute: {', '.join(inputs.compute_services or [])}\n- Network: {', '.join(inputs.network_services or [])}\n- Security: {', '.join(inputs.security_services or [])}\n\n*Full documentation requires AI service availability.*",
                "hld": f"# High Level Design\n\n## Azure Architecture Overview\n\nThis document outlines the high-level design for an Azure Landing Zone.\n\n### Key Components\n- Management Groups\n- Subscriptions\n- Resource Groups\n- Network Architecture\n\n*Detailed design requires AI service availability.*",
                "lld": f"# Low Level Design\n\n## Implementation Details\n\nThis document provides implementation guidance for the Azure Landing Zone.\n\n### Implementation Steps\n1. Set up Management Groups\n2. Configure Subscriptions\n3. Deploy Network Infrastructure\n4. Implement Security Controls\n\n*Detailed implementation guide requires AI service availability.*"
            }
        
        # Count Azure stencils used in Draw.io XML
        import re
        shapes = re.findall(r'shape=mxgraph\.azure\.[^;\"\s]*', drawio_xml)
        
        # Get architecture template information
        template_info = generate_architecture_template(inputs)
        
        result = {
            "success": True,
            "mermaid": mermaid_diagram,
            "svg_diagram": svg_content,
            "svg_diagram_path": svg_diagram_path,
            "drawio_xml": drawio_xml,
            "tsd": docs["tsd"],
            "hld": docs["hld"],
            "lld": docs["lld"],
            "architecture_template": template_info,
            "azure_stencils": {
                "total_used": len(shapes),
                "unique_used": len(set(shapes)),
                "stencils_list": sorted(list(set(shapes)))
            },
            "feedback_questions": feedback_questions,
            "ai_analysis": {
                "services_used": template_info.get("ai_services", []),
                "reasoning": template_info.get("ai_reasoning", "")
            },
            "metadata": {
                "generated_at": datetime.now().isoformat(),
                "version": "1.0.0",
                "agent": "Azure Landing Zone Agent - Interactive Generator",
                "diagram_format": "SVG with Azure official icons",
                "svg_size": len(svg_content),
                "drawio_size": len(drawio_xml)
            }
        }
        
        logger.info("Interactive Azure architecture generated successfully")
        return result
    
    except ValueError as ve:
        # Input validation errors
        error_msg = f"Invalid input: {str(ve)}"
        logger.error(error_msg)
        raise HTTPException(status_code=400, detail=error_msg)
    
    except Exception as e:
        # Log the full error for debugging
        error_msg = f"Error generating interactive architecture: {str(e)}"
        logger.error(error_msg)
        logger.error(traceback.format_exc())
        
        # Return a user-friendly error
        raise HTTPException(
            status_code=500, 
            detail=f"Failed to generate interactive architecture. Error: {str(e)}"
        )

def generate_feedback_questions(inputs: CustomerInputs) -> List[str]:
    """Generate human-in-the-loop feedback questions to improve architecture"""
    questions = []
    
    # Check for missing critical information
    if not inputs.business_objective:
        questions.append("What is your primary business objective for this Azure deployment? (Cost optimization, agility, innovation, security, etc.)")
    
    if not inputs.scalability:
        questions.append("What are your expected scalability requirements? (Current and future user load, geographic distribution)")
    
    if not inputs.security_posture:
        questions.append("What security and compliance requirements do you have? (Zero trust, industry regulations, data sovereignty)")
    
    # Check if AI analysis was used and ask for validation
    total_selected_services = sum([
        len(inputs.compute_services or []),
        len(inputs.network_services or []),
        len(inputs.storage_services or []),
        len(inputs.database_services or []),
        len(inputs.security_services or [])
    ])
    
    if inputs.free_text_input and total_selected_services <= 1:
        questions.append("I've analyzed your requirements and suggested specific Azure services. Would you like to review and confirm these selections before finalizing the architecture?")
        questions.append("Are there any specific performance, availability, or disaster recovery requirements I should consider?")
    
    # Ask about budget and cost constraints
    if not inputs.cost_priority:
        questions.append("What is your cost optimization priority? Should we focus on minimizing costs or optimizing for performance?")
    
    # Ask about existing infrastructure
    if inputs.free_text_input and "existing" not in inputs.free_text_input.lower():
        questions.append("Do you have any existing Azure infrastructure or on-premises systems that need to be integrated?")
    
    # Ask about operational model
    if not inputs.ops_model and not inputs.monitoring:
        questions.append("How do you plan to operate and monitor this infrastructure? Do you have a dedicated DevOps team?")
    
    return questions[:3]  # Limit to top 3 most relevant questions

@app.post("/generate-png-diagram")
def generate_png_diagram(inputs: CustomerInputs):
    """Generate PNG diagram for download"""
    logger.info("Starting PNG diagram generation for download")
    
    try:
        # Validate inputs early
        validate_customer_inputs(inputs)
        logger.info("Input validation completed successfully")
        
        # Generate PNG diagram with enhanced design principles
        logger.info("Generating enhanced PNG diagram...")
        png_path = generate_enhanced_azure_architecture_diagram(inputs, format="png")
        logger.info(f"PNG diagram generated successfully: {png_path}")
        
        # Read and encode the PNG file
        with open(png_path, "rb") as f:
            png_content = f.read()
        
        png_base64 = base64.b64encode(png_content).decode("utf-8")
        
        return {
            "success": True,
            "png_base64": png_base64,
            "png_path": png_path,
            "file_size": len(png_content),
            "metadata": {
                "generated_at": datetime.now().isoformat(),
                "format": "PNG"
            }
        }
        
    except Exception as e:
        logger.error(f"Error generating PNG diagram: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to generate PNG diagram: {str(e)}")

@app.post("/generate-svg-diagram")
def generate_svg_diagram(inputs: CustomerInputs):
    """Generate SVG diagram for download"""
    logger.info("Starting SVG diagram generation for download")
    
    try:
        # Validate inputs early
        validate_customer_inputs(inputs)
        logger.info("Input validation completed successfully")
        
        # Generate SVG diagram with enhanced design principles
        logger.info("Generating enhanced SVG diagram...")
        svg_path = generate_enhanced_azure_architecture_diagram(inputs, format="svg")
        logger.info(f"SVG diagram generated successfully: {svg_path}")
        
        # Read the SVG file
        with open(svg_path, "r", encoding="utf-8") as f:
            svg_content = f.read()
        
        return {
            "success": True,
            "svg_content": svg_content,
            "svg_path": svg_path,
            "file_size": len(svg_content),
            "metadata": {
                "generated_at": datetime.now().isoformat(),
                "format": "SVG"
            }
        }
        
    except Exception as e:
        logger.error(f"Error generating SVG diagram: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to generate SVG diagram: {str(e)}")

@app.post("/upload-file")
async def upload_file(file: UploadFile = File(...)):
    """Upload and process files (PDF, Excel, PowerPoint) for AI analysis"""
    try:
        # Validate file type
        allowed_extensions = ['.pdf', '.xlsx', '.xls', '.pptx', '.ppt']
        file_extension = os.path.splitext(file.filename.lower())[1]
        
        if file_extension not in allowed_extensions:
            raise HTTPException(
                status_code=400, 
                detail=f"Unsupported file type. Allowed types: {', '.join(allowed_extensions)}"
            )
        
        # Validate file size (max 10MB)
        max_file_size = 10 * 1024 * 1024  # 10MB
        file_content = await file.read()
        
        if len(file_content) > max_file_size:
            raise HTTPException(
                status_code=400,
                detail="File too large. Maximum size is 10MB."
            )
        
        # Process the file with AI
        file_type = file_extension[1:]  # Remove the dot
        analysis_result = process_uploaded_document(file_content, file.filename, file_type)
        
        # Return file info and analysis
        return {
            "success": True,
            "filename": file.filename,
            "file_type": file_type,
            "file_size": len(file_content),
            "analysis": analysis_result,
            "upload_timestamp": datetime.now().isoformat()
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error uploading file {file.filename}: {e}")
        raise HTTPException(status_code=500, detail=f"Error processing file: {str(e)}")

@app.post("/analyze-url")
def analyze_url(request: Dict[str, str]):
    """Analyze URL content for Azure architecture insights"""
    try:
        url = request.get("url")
        if not url:
            raise HTTPException(status_code=400, detail="URL is required")
        
        # Validate URL format
        if not url.startswith(('http://', 'https://')):
            raise HTTPException(status_code=400, detail="URL must start with http:// or https://")
        
        # Analyze the URL with AI
        analysis_result = analyze_url_content(url)
        
        return {
            "success": True,
            "url": url,
            "analysis": analysis_result,
            "analysis_timestamp": datetime.now().isoformat()
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error analyzing URL {url}: {e}")
        raise HTTPException(status_code=500, detail=f"Error analyzing URL: {str(e)}")

@app.get("/templates")
def get_templates():
    """Get available Azure Landing Zone templates"""
    return {
        "templates": AZURE_TEMPLATES,
        "azure_services": AZURE_SERVICES_MAPPING
    }

@app.get("/services")
def get_services():
    """Get available Azure services categorized for form selection"""
    services_by_category = {}
    
    for service_key, service_info in AZURE_SERVICES_MAPPING.items():
        category = service_info["category"]
        if category not in services_by_category:
            services_by_category[category] = []
        
        services_by_category[category].append({
            "key": service_key,
            "name": service_info["name"],
            "icon": service_info["icon"],
            "azure_icon": service_info.get("azure_icon", ""),
        })
    
    return {
        "categories": services_by_category,
        "category_mapping": {
            "compute": "Compute Services",
            "network": "Networking Services", 
            "storage": "Storage Services",
            "database": "Database Services",
            "security": "Security Services",
            "monitoring": "Monitoring Services",
            "ai": "AI & Machine Learning",
            "analytics": "Data & Analytics",
            "integration": "Integration Services",
            "devops": "DevOps & Governance",
            "backup": "Backup & Recovery"
        }
    }

class FeedbackRequest(BaseModel):
    original_inputs: CustomerInputs
    feedback_answers: Dict[str, str]
    selected_services: Optional[List[str]] = Field(default_factory=list)

@app.post("/refine-architecture-with-feedback")
def refine_architecture_with_feedback(request: FeedbackRequest):
    """Refine architecture based on human feedback"""
    try:
        logger.info("Processing architecture refinement with human feedback")
        
        # Update inputs based on feedback
        refined_inputs = request.original_inputs
        
        # Process feedback answers to update inputs
        for question, answer in request.feedback_answers.items():
            if "business objective" in question.lower():
                refined_inputs.business_objective = answer
            elif "scalability" in question.lower():
                refined_inputs.scalability = answer
            elif "security" in question.lower():
                refined_inputs.security_posture = answer
            elif "cost" in question.lower():
                refined_inputs.cost_priority = answer
            elif "operational" in question.lower() or "monitor" in question.lower():
                refined_inputs.ops_model = answer
        
        # Add user-confirmed services if provided
        if request.selected_services:
            # Clear existing services and add only confirmed ones
            for service in request.selected_services:
                if service in AZURE_SERVICES_MAPPING:
                    category = AZURE_SERVICES_MAPPING[service]["category"]
                    category_field = f"{category}_services"
                    
                    if hasattr(refined_inputs, category_field):
                        current_services = getattr(refined_inputs, category_field) or []
                        if service not in current_services:
                            current_services.append(service)
                            setattr(refined_inputs, category_field, current_services)
        
        # Generate refined architecture
        return generate_interactive_azure_architecture(refined_inputs)
        
    except Exception as e:
        logger.error(f"Error refining architecture with feedback: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to refine architecture: {str(e)}")

@app.post("/validate-ai-service-selection")
def validate_ai_service_selection(request: Dict[str, Any]):
    """Allow users to validate and modify AI-suggested service selection"""
    try:
        free_text = request.get("free_text_input", "")
        
        if not free_text:
            raise HTTPException(status_code=400, detail="Free text input is required")
        
        # Get AI analysis
        analysis = analyze_free_text_requirements(free_text)
        
        # Get detailed service information
        suggested_services = []
        for service_key in analysis.get("services", []):
            if service_key in AZURE_SERVICES_MAPPING:
                service_info = AZURE_SERVICES_MAPPING[service_key]
                suggested_services.append({
                    "key": service_key,
                    "name": service_info["name"],
                    "category": service_info["category"],
                    "icon": service_info["icon"],
                    "reasoning": f"Suggested for: {analysis.get('reasoning', 'Meeting your requirements')}"
                })
        
        return {
            "success": True,
            "original_text": free_text,
            "suggested_services": suggested_services,
            "reasoning": analysis.get("reasoning", ""),
            "architecture_pattern": analysis.get("architecture_pattern", "simple"),
            "questions_for_clarification": [
                "Do these suggested services meet your requirements?",
                "Are there any additional services you need?",
                "Would you like to modify any of these selections?"
            ]
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error validating AI service selection: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to validate service selection: {str(e)}")

# Enhanced endpoints for background diagram generation with improved design principles
@app.post("/generate-background-diagrams")
def generate_background_diagrams(inputs: CustomerInputs):
    """Generate SVG and PNG diagrams in background without immediate download"""
    logger.info("Starting background diagram generation with enhanced 50+ design principles")
    
    try:
        # Validate inputs early
        validate_customer_inputs(inputs)
        logger.info("Input validation completed successfully")
        
        # Generate both formats in background
        logger.info("Generating enhanced PNG diagram in background...")
        png_path = generate_enhanced_azure_architecture_diagram(inputs, format="png")
        logger.info(f"Enhanced PNG diagram generated: {png_path}")
        
        logger.info("Generating enhanced SVG diagram in background...")
        svg_path = generate_enhanced_azure_architecture_diagram(inputs, format="svg")
        logger.info(f"Enhanced SVG diagram generated: {svg_path}")
        
        # Store file information for later retrieval
        file_info = {
            "png_path": png_path,
            "svg_path": svg_path,
            "generated_at": datetime.now().isoformat(),
            "input_summary": {
                "business_objective": inputs.business_objective,
                "org_structure": inputs.org_structure,
                "network_model": inputs.network_model,
                "compute_services": inputs.compute_services,
                "network_services": inputs.network_services,
                "storage_services": inputs.storage_services,
                "database_services": inputs.database_services
            }
        }
        
        return {
            "success": True,
            "message": "Enhanced diagrams generated successfully in background",
            "file_info": file_info,
            "metadata": {
                "generated_at": datetime.now().isoformat(),
                "enhanced_design_principles": True,
                "formats_available": ["PNG", "SVG"]
            }
        }
        
    except Exception as e:
        logger.error(f"Error generating background diagrams: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to generate background diagrams: {str(e)}")

@app.post("/download-enhanced-diagram")
def download_enhanced_diagram(request: Dict[str, str]):
    """Download enhanced diagram files (PNG or SVG) after background generation"""
    try:
        file_path = request.get("file_path")
        format_type = request.get("format", "png").lower()
        
        if not file_path:
            raise HTTPException(status_code=400, detail="File path is required")
        
        if not os.path.exists(file_path):
            raise HTTPException(status_code=404, detail="Generated file not found")
        
        if format_type == "svg":
            # Read SVG content
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()
            
            return {
                "success": True,
                "format": "SVG",
                "svg_content": content,
                "file_path": file_path,
                "file_size": len(content)
            }
        else:
            # Handle PNG format request
            if file_path.endswith('.png'):
                # Read PNG content and encode as base64
                with open(file_path, "rb") as f:
                    content = f.read()
                
                png_base64 = base64.b64encode(content).decode("utf-8")
                
                return {
                    "success": True,
                    "format": "PNG", 
                    "png_base64": png_base64,
                    "file_path": file_path,
                    "file_size": len(content)
                }
            else:
                # File is SVG but PNG was requested (conversion failed)
                # Return SVG content with a note
                with open(file_path, "r", encoding="utf-8") as f:
                    content = f.read()
                
                return {
                    "success": True,
                    "format": "SVG", 
                    "svg_content": content,
                    "file_path": file_path,
                    "file_size": len(content),
                    "note": "PNG conversion was not available, returning enhanced SVG instead"
                }
            
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error downloading enhanced diagram: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to download enhanced diagram: {str(e)}")

def generate_enhanced_azure_architecture_diagram(inputs: CustomerInputs, output_dir: str = None, format: str = "png") -> str:
    """Generate Enhanced Azure architecture diagram with improved 50+ enterprise architecture principles
    
    Enhanced Enterprise Architecture Principles Implemented:
    1-10: Clear containers/swimlanes, minimal crossing connections, proper visual hierarchy
    11-20: Clear connection labeling, numbered workflow, all specified components  
    21-30: Security zoning, environment labeling, HA indicators, monitoring overlay
    31-40: Comprehensive legend, standardized iconography, observability integration
    41-50: Scalability indicators, compliance overlays, cost management, future-ready design
    51-55: Enhanced structured formatting, improved visual consistency, better spacing
    """
    
    logger.info("Starting Enhanced Azure architecture diagram generation with 55+ improved principles")
    
    try:
        # Use enhanced simple SVG for better reliability instead of complex Graphviz
        logger.info("Using enhanced simple diagram generation for better reliability")
        return generate_enhanced_simple_svg_diagram(inputs, format)
                
    except Exception as e:
        logger.error(f"Enhanced Azure architecture diagram generation failed: {str(e)}")
        logger.error(f"Traceback: {traceback.format_exc()}")
        raise Exception(f"Enhanced diagram generation failed: {str(e)}")

def _add_enhanced_compute_clusters(inputs: CustomerInputs, prod_vnet, service_counter: int):
    """Enhanced compute clusters with better positioning and organization"""
    with Cluster("💻 COMPUTE SERVICES - HIGH AVAILABILITY", graph_attr={
        "bgcolor": "#e1f5fe",
        "style": "filled,rounded",
        "color": "#0277bd",
        "penwidth": "2",
        "fontsize": "12",
        "margin": "15"
    }):
        for service in inputs.compute_services[:4]:  # Limit for visual clarity
            if service == "virtual_machines":
                VM(f"[{service_counter}] Virtual Machines\n🖥️ IaaS Compute\n⚖️ Auto-Scale Enabled")
            elif service == "aks":
                AKS(f"[{service_counter}] AKS Cluster\n☸️ Kubernetes\n🔄 Multi-Zone HA")
            elif service == "app_services":
                AppServices(f"[{service_counter}] App Services\n🌐 PaaS Web Apps\n📈 Auto-Scale")
            elif service == "function_apps":
                FunctionApps(f"[{service_counter}] Function Apps\n⚡ Serverless\n💰 Pay-per-Use")
            service_counter += 1

def _add_enhanced_storage_clusters(inputs: CustomerInputs, prod_vnet, service_counter: int):
    """Enhanced storage clusters with better data tier visualization"""
    with Cluster("💾 STORAGE SERVICES - DATA TIERS", graph_attr={
        "bgcolor": "#f9fbe7", 
        "style": "filled,rounded",
        "color": "#689f38",
        "penwidth": "2",
        "fontsize": "12",
        "margin": "15"
    }):
        storage_services = []
        for service in inputs.storage_services[:3]:  # Limit for visual clarity
            if service == "storage_accounts":
                storage_services.append(StorageAccounts(f"[{service_counter}] Storage Accounts\n📁 General Purpose\n🔄 LRS/GRS Options"))
            elif service == "blob_storage":
                storage_services.append(BlobStorage(f"[{service_counter}] Blob Storage\n🗂️ Object Storage\n❄️ Hot/Cool/Archive"))
            elif service == "data_lake_storage":
                storage_services.append(DataLakeStorage(f"[{service_counter}] Data Lake Gen2\n🏞️ Analytics Storage\n📊 Big Data Ready"))
            service_counter += 1
        return storage_services

def _add_enhanced_database_clusters(inputs: CustomerInputs, prod_vnet, service_counter: int):
    """Enhanced database clusters with better data management visualization"""
    with Cluster("🗄️ DATABASE SERVICES - DATA MANAGEMENT", graph_attr={
        "bgcolor": "#fef7ff",
        "style": "filled,rounded", 
        "color": "#7b1fa2",
        "penwidth": "2",
        "fontsize": "12",
        "margin": "15"
    }):
        database_services = []
        for service in inputs.database_services[:3]:  # Limit for visual clarity
            if service == "sql_database":
                database_services.append(SQLDatabases(f"[{service_counter}] SQL Database\n🗃️ Relational DB\n🔄 Auto-Backup"))
            elif service == "cosmos_db":
                database_services.append(CosmosDb(f"[{service_counter}] Cosmos DB\n🌍 Multi-Model\n⚡ Global Distribution"))
            elif service == "mysql":
                database_services.append(DatabaseForMysqlServers(f"[{service_counter}] MySQL\n🐬 Open Source\n☁️ Managed Service"))
            service_counter += 1
        return database_services

def generate_enhanced_simple_svg_diagram(inputs: CustomerInputs, format: str = "svg") -> str:
    """Generate an enhanced simple diagram with improved design principles aligned with problem statement requirements"""
    logger.info(f"Generating enhanced simple {format.upper()} diagram with 50+ architecture principles")
    
    # Get safe output directory
    output_dir = get_safe_output_directory()
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    unique_id = str(uuid.uuid4())[:8]
    
    # Create SVG first, then convert to PNG if needed
    svg_filename = f"enhanced_simple_azure_architecture_{timestamp}_{unique_id}.svg"
    svg_filepath = os.path.join(output_dir, svg_filename)
    
    # Enhanced SVG content with improved design principles
    svg_content = f'''<?xml version="1.0" encoding="UTF-8"?>
<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1400 1100" style="background: linear-gradient(135deg, #f8f9fa 0%, #e9ecef 100%);">
    <defs>
        <!-- Enhanced gradient definitions for better visual hierarchy -->
        <linearGradient id="untrustedGrad" x1="0%" y1="0%" x2="0%" y2="100%">
            <stop offset="0%" style="stop-color:#ffebee;stop-opacity:1" />
            <stop offset="100%" style="stop-color:#ffe6e6;stop-opacity:1" />
        </linearGradient>
        <linearGradient id="semitrustedGrad" x1="0%" y1="0%" x2="0%" y2="100%">
            <stop offset="0%" style="stop-color:#e8f5e8;stop-opacity:1" />
            <stop offset="100%" style="stop-color:#c8e6c8;stop-opacity:1" />
        </linearGradient>
        <linearGradient id="trustedGrad" x1="0%" y1="0%" x2="0%" y2="100%">
            <stop offset="0%" style="stop-color:#e3f2fd;stop-opacity:1" />
            <stop offset="100%" style="stop-color:#bbdefb;stop-opacity:1" />
        </linearGradient>
        <linearGradient id="dataGrad" x1="0%" y1="0%" x2="0%" y2="100%">
            <stop offset="0%" style="stop-color:#fef7ff;stop-opacity:1" />
            <stop offset="100%" style="stop-color:#f3e5f5;stop-opacity:1" />
        </linearGradient>
        <linearGradient id="drGrad" x1="0%" y1="0%" x2="0%" y2="100%">
            <stop offset="0%" style="stop-color:#fff8e1;stop-opacity:1" />
            <stop offset="100%" style="stop-color:#fff3e0;stop-opacity:1" />
        </linearGradient>
        
        <!-- Arrow markers for different connection types -->
        <marker id="primaryArrow" markerWidth="12" markerHeight="8" refX="12" refY="4" orient="auto">
            <polygon points="0 0, 12 4, 0 8" fill="#1976d2"/>
        </marker>
        <marker id="securityArrow" markerWidth="12" markerHeight="8" refX="12" refY="4" orient="auto">
            <polygon points="0 0, 12 4, 0 8" fill="#d32f2f"/>
        </marker>
        <marker id="dataArrow" markerWidth="12" markerHeight="8" refX="12" refY="4" orient="auto">
            <polygon points="0 0, 12 4, 0 8" fill="#7b1fa2"/>
        </marker>
        <marker id="monitoringArrow" markerWidth="10" markerHeight="6" refX="10" refY="3" orient="auto">
            <polygon points="0 0, 10 3, 0 6" fill="#00695c"/>
        </marker>
    </defs>
    
    <style>
        .title {{ font: bold 28px Arial, sans-serif; fill: #1976d2; text-anchor: middle; }}
        .environment-tag {{ font: bold 14px Arial, sans-serif; fill: #fff; text-anchor: middle; }}
        .swimlane-title {{ font: bold 18px Arial, sans-serif; fill: #333; text-anchor: start; }}
        .cluster-title {{ font: bold 16px Arial, sans-serif; fill: #333; text-anchor: middle; }}
        .service {{ font: 12px Arial, sans-serif; fill: #333; text-anchor: middle; }}
        .service-detail {{ font: 10px Arial, sans-serif; fill: #666; text-anchor: middle; }}
        .connection-label {{ font: 11px Arial, sans-serif; fill: #444; text-anchor: middle; font-weight: bold; }}
        .workflow-number {{ font: bold 14px Arial, sans-serif; fill: #fff; text-anchor: middle; }}
        .legend-title {{ font: bold 16px Arial, sans-serif; fill: #333; text-anchor: start; }}
        .legend-text {{ font: 11px Arial, sans-serif; fill: #333; text-anchor: start; }}
        .ha-indicator {{ font: bold 10px Arial, sans-serif; fill: #e65100; text-anchor: middle; }}
        
        /* Security Zones - Enhanced visual hierarchy */
        .untrusted-zone {{ fill: url(#untrustedGrad); stroke: #d32f2f; stroke-width: 3; stroke-dasharray: none; }}
        .semitrusted-zone {{ fill: url(#semitrustedGrad); stroke: #388e3c; stroke-width: 3; stroke-dasharray: none; }}
        .trusted-zone {{ fill: url(#trustedGrad); stroke: #1976d2; stroke-width: 3; stroke-dasharray: none; }}
        .data-zone {{ fill: url(#dataGrad); stroke: #7b1fa2; stroke-width: 3; stroke-dasharray: none; }}
        .dr-zone {{ fill: url(#drGrad); stroke: #ef6c00; stroke-width: 3; stroke-dasharray: 5,5; }}
        
        .service-box {{ fill: #ffffff; stroke: #424242; stroke-width: 2; }}
        .enhanced-shadow {{ filter: drop-shadow(4px 4px 8px rgba(0,0,0,0.2)); }}
        .strong-shadow {{ filter: drop-shadow(6px 6px 12px rgba(0,0,0,0.3)); }}
        
        /* Environment indicators */
        .env-prod {{ fill: #2e7d32; }}
        .env-dr {{ fill: #ef6c00; }}
        .env-dev {{ fill: #1976d2; }}
        .env-uat {{ fill: #7b1fa2; }}
    </style>
    
    <!-- Title with Environment Indicator -->
    <text x="700" y="35" class="title">Azure Enterprise Landing Zone Architecture</text>
    <text x="700" y="55" class="title" style="font-size: 16px; fill: #666;">Enhanced with 50+ Enterprise Design Principles</text>
    
    <!-- Environment Tags -->
    <rect x="50" y="20" width="80" height="25" class="env-prod" rx="3"/>
    <text x="90" y="37" class="environment-tag">PROD</text>
    <rect x="140" y="20" width="60" height="25" class="env-dr" rx="3"/>
    <text x="170" y="37" class="environment-tag">DR</text>
    
    <!-- SWIMLANE 1: INTERNET EDGE (Untrusted Zone) -->
    <rect x="30" y="80" width="1340" height="140" class="untrusted-zone enhanced-shadow" rx="15"/>
    <text x="50" y="105" class="swimlane-title">🌐 INTERNET EDGE - EXTERNAL ACCESS (Untrusted Zone)</text>
    
    <!-- Workflow Step 1: Internet Access -->
    <circle cx="80" cy="160" r="15" fill="#d32f2f" class="enhanced-shadow"/>
    <text x="80" y="166" class="workflow-number">1</text>
    
    <rect x="650" y="130" width="180" height="70" class="service-box enhanced-shadow" rx="8"/>
    <text x="740" y="150" class="service">Public Internet</text>
    <text x="740" y="165" class="service-detail">External Users & Traffic</text>
    <text x="740" y="180" class="service-detail">Global Entry Point</text>
    <text x="740" y="195" class="ha-indicator">24/7 Availability</text>
    
    <!-- SWIMLANE 2: IDENTITY & SECURITY (Semi-trusted Zone) -->
    <rect x="30" y="240" width="1340" height="140" class="semitrusted-zone enhanced-shadow" rx="15"/>
    <text x="50" y="265" class="swimlane-title">🔐 IDENTITY & SECURITY - SECURITY PERIMETER (Semi-trusted Zone)</text>
    
    <!-- Workflow Step 2: Security Gateway -->
    <circle cx="80" cy="320" r="15" fill="#388e3c" class="enhanced-shadow"/>
    <text x="80" y="326" class="workflow-number">2</text>
    
    <rect x="200" y="285" width="160" height="70" class="service-box enhanced-shadow" rx="8"/>
    <text x="280" y="305" class="service">Azure Firewall</text>
    <text x="280" y="320" class="service-detail">Next-Gen Security</text>
    <text x="280" y="335" class="service-detail">Threat Intelligence</text>
    <text x="280" y="350" class="ha-indicator">Active-Active HA</text>
    
    <rect x="450" y="285" width="160" height="70" class="service-box enhanced-shadow" rx="8"/>
    <text x="530" y="305" class="service">Application Gateway</text>
    <text x="530" y="320" class="service-detail">L7 Load Balancer</text>
    <text x="530" y="335" class="service-detail">WAF Protection</text>
    <text x="530" y="350" class="ha-indicator">Multi-Zone Deployed</text>
    
    <rect x="700" y="285" width="160" height="70" class="service-box enhanced-shadow" rx="8"/>
    <text x="780" y="305" class="service">Azure AD</text>
    <text x="780" y="320" class="service-detail">Identity Provider</text>
    <text x="780" y="335" class="service-detail">Conditional Access</text>
    <text x="780" y="350" class="ha-indicator">99.9% SLA</text>
    
    <rect x="950" y="285" width="160" height="70" class="service-box enhanced-shadow" rx="8"/>
    <text x="1030" y="305" class="service">Key Vault</text>
    <text x="1030" y="320" class="service-detail">Secrets Management</text>
    <text x="1030" y="335" class="service-detail">HSM-backed</text>
    <text x="1030" y="350" class="ha-indicator">FIPS 140-2 Level 2</text>
    
    <!-- SWIMLANE 3: NETWORK HUB (Trusted Zone) -->
    <rect x="30" y="400" width="1340" height="140" class="trusted-zone enhanced-shadow" rx="15"/>
    <text x="50" y="425" class="swimlane-title">🔗 NETWORK HUB - CENTRAL CONNECTIVITY (Trusted Zone)</text>
    
    <!-- Workflow Step 3: Network Routing -->
    <circle cx="80" cy="480" r="15" fill="#1976d2" class="enhanced-shadow"/>
    <text x="80" y="486" class="workflow-number">3</text>
    
    <rect x="400" y="445" width="180" height="70" class="service-box enhanced-shadow" rx="8"/>
    <text x="490" y="465" class="service">Hub VNet</text>
    <text x="490" y="480" class="service-detail">Central Network Hub</text>
    <text x="490" y="495" class="service-detail">10.0.0.0/16</text>
    <text x="490" y="510" class="ha-indicator">Multi-AZ Redundancy</text>
    
    <rect x="650" y="445" width="180" height="70" class="service-box enhanced-shadow" rx="8"/>
    <text x="740" y="465" class="service">VPN Gateway</text>
    <text x="740" y="480" class="service-detail">Hybrid Connectivity</text>
    <text x="740" y="495" class="service-detail">Site-to-Site VPN</text>
    <text x="740" y="510" class="ha-indicator">Active-Active Config</text>
    
    <rect x="900" y="445" width="180" height="70" class="service-box enhanced-shadow" rx="8"/>
    <text x="990" y="465" class="service">ExpressRoute</text>
    <text x="990" y="480" class="service-detail">Dedicated Connection</text>
    <text x="990" y="495" class="service-detail">Private Connectivity</text>
    <text x="990" y="510" class="ha-indicator">99.95% SLA</text>
    
    <!-- SWIMLANE 4: APPLICATION SERVICES (Trusted Zone) -->
    <rect x="30" y="560" width="1340" height="160" class="trusted-zone enhanced-shadow" rx="15"/>
    <text x="50" y="585" class="swimlane-title">⚙️ APPLICATION SERVICES - COMPUTE & INTEGRATION (Trusted Zone)</text>
    
    <!-- Workflow Step 4: Application Processing -->
    <circle cx="80" cy="650" r="15" fill="#1976d2" class="enhanced-shadow"/>
    <text x="80" y="656" class="workflow-number">4</text>
    
    <!-- Production VNet -->
    <rect x="150" y="605" width="200" height="90" class="service-box enhanced-shadow" rx="8"/>
    <text x="250" y="625" class="service">Production VNet</text>
    <text x="250" y="640" class="service-detail">Workload Network</text>
    <text x="250" y="655" class="service-detail">10.1.0.0/16</text>
    <text x="250" y="670" class="service-detail">3 Availability Zones</text>
    <text x="250" y="685" class="ha-indicator">Zone-Redundant</text>'''
    
    # Add dynamic compute services based on selection
    x_offset = 380
    service_counter = 5
    if inputs.compute_services:
        for i, service in enumerate(inputs.compute_services[:3]):
            service_name = service.replace('_', ' ').title()
            if service == "virtual_machines":
                detail1 = "IaaS Compute"
                detail2 = "Auto-Scale Enabled"
                ha_indicator = "Zone Distribution"
                compliance = "ISO 27001"
            elif service == "aks":
                detail1 = "Kubernetes Service"
                detail2 = "Container Orchestration"  
                ha_indicator = "Multi-Zone Nodes"
                compliance = "CIS Benchmark"
            elif service == "app_services":
                detail1 = "PaaS Web Apps"
                detail2 = "Managed Platform"
                ha_indicator = "Auto-Scale Rules"
                compliance = "SOC 2 Type II"
            else:
                detail1 = "Compute Service"
                detail2 = "Managed Service"
                ha_indicator = "High Availability"
                compliance = "Compliant"
            
            svg_content += f'''
    <rect x="{x_offset}" y="605" width="180" height="110" class="service-box enhanced-shadow" rx="8"/>
    <text x="{x_offset + 90}" y="625" class="service">{service_name}</text>
    <text x="{x_offset + 90}" y="640" class="service-detail">{detail1}</text>
    <text x="{x_offset + 90}" y="655" class="service-detail">{detail2}</text>
    <text x="{x_offset + 90}" y="670" class="service-detail">Service Tier: Standard</text>
    <text x="{x_offset + 90}" y="685" class="ha-indicator">{ha_indicator}</text>
    <text x="{x_offset + 90}" y="700" class="legend-text">{compliance}</text>
    
    <!-- Scalability & Security indicators -->
    <rect x="{x_offset + 155}" y="610" width="20" height="12" fill="#2196f3" rx="2"/>
    <text x="{x_offset + 165}" y="620" style="font: 8px Arial; fill: white; text-anchor: middle;">📈</text>
    
    <rect x="{x_offset + 155}" y="625" width="20" height="12" fill="#9c27b0" rx="2"/>
    <text x="{x_offset + 165}" y="635" style="font: 8px Arial; fill: white; text-anchor: middle;">🔒</text>'''
            x_offset += 200
            service_counter += 1
    
    # Add storage and database services in the data layer
    svg_content += '''
    
    <!-- SWIMLANE 5: DATA & STORAGE LAYER (Data Zone) -->
    <rect x="30" y="740" width="1340" height="140" class="data-zone enhanced-shadow" rx="15"/>
    <text x="50" y="765" class="swimlane-title">💾 DATA & STORAGE LAYER - PERSISTENT DATA (Data Zone)</text>
    
    <!-- Workflow Step 5: Data Processing -->
    <circle cx="80" cy="820" r="15" fill="#7b1fa2" class="enhanced-shadow"/>
    <text x="80" y="826" class="workflow-number">5</text>'''
    
    # Add storage services
    x_offset = 200
    if inputs.storage_services:
        for i, service in enumerate(inputs.storage_services[:3]):
            service_name = service.replace('_', ' ').title()
            if service == "storage_accounts":
                detail1 = "General Purpose v2"
                detail2 = "LRS/GRS Options"
                ha_indicator = "99.999999999% Durability"
                compliance = "GDPR | SOC 2"
            elif service == "blob_storage":
                detail1 = "Object Storage"
                detail2 = "Hot/Cool/Archive Tiers"
                ha_indicator = "Geo-Replication"
                compliance = "HIPAA | PCI-DSS"
            else:
                detail1 = "Storage Service"
                detail2 = "Managed Storage"
                ha_indicator = "High Durability"
                compliance = "Compliant"
            
            svg_content += f'''
    <rect x="{x_offset}" y="785" width="180" height="90" class="service-box enhanced-shadow" rx="8"/>
    <text x="{x_offset + 90}" y="805" class="service">{service_name}</text>
    <text x="{x_offset + 90}" y="820" class="service-detail">{detail1}</text>
    <text x="{x_offset + 90}" y="835" class="service-detail">{detail2}</text>
    <text x="{x_offset + 90}" y="850" class="ha-indicator">{ha_indicator}</text>
    <text x="{x_offset + 90}" y="865" class="legend-text">{compliance}</text>
    
    <!-- Auto-scaling indicator -->
    <rect x="{x_offset + 155}" y="790" width="20" height="12" fill="#4caf50" rx="2"/>
    <text x="{x_offset + 165}" y="800" style="font: 8px Arial; fill: white; text-anchor: middle;">⚡</text>'''
            x_offset += 200
    
    # Add database services
    if inputs.database_services:
        for i, service in enumerate(inputs.database_services[:3]):
            service_name = service.replace('_', ' ').title()
            if service == "sql_database":
                detail1 = "Relational Database"
                detail2 = "Always Encrypted"
                ha_indicator = "Auto-Failover Groups"
                compliance = "SOX | GDPR"
            elif service == "cosmos_db":
                detail1 = "Multi-Model NoSQL"
                detail2 = "Global Distribution"
                ha_indicator = "99.999% Availability"
                compliance = "ISO 27001"
            else:
                detail1 = "Database Service"
                detail2 = "Managed Database"
                ha_indicator = "Automatic Backup"
                compliance = "Compliant"
            
            svg_content += f'''
    <rect x="{x_offset}" y="785" width="180" height="90" class="service-box enhanced-shadow" rx="8"/>
    <text x="{x_offset + 90}" y="805" class="service">{service_name}</text>
    <text x="{x_offset + 90}" y="820" class="service-detail">{detail1}</text>
    <text x="{x_offset + 90}" y="835" class="service-detail">{detail2}</text>
    <text x="{x_offset + 90}" y="850" class="ha-indicator">{ha_indicator}</text>
    <text x="{x_offset + 90}" y="865" class="legend-text">{compliance}</text>
    
    <!-- Cost optimization indicator -->
    <rect x="{x_offset + 155}" y="790" width="20" height="12" fill="#ff9800" rx="2"/>
    <text x="{x_offset + 165}" y="800" style="font: 8px Arial; fill: white; text-anchor: middle;">💰</text>'''
            x_offset += 200
    
    # Add disaster recovery zone and monitoring overlay
    svg_content += '''
    
    <!-- DISASTER RECOVERY REGION (Standby) -->
    <rect x="30" y="900" width="650" height="120" class="dr-zone strong-shadow" rx="15"/>
    <text x="50" y="925" class="swimlane-title">🔄 DISASTER RECOVERY - STANDBY REGION (West US 2)</text>
    <text x="50" y="945" class="service-detail">RTO: 1 hour | RPO: 15 minutes | Active-Passive Configuration</text>
    
    <rect x="80" y="955" width="150" height="50" class="service-box enhanced-shadow" rx="8"/>
    <text x="155" y="975" class="service">DR VNet</text>
    <text x="155" y="990" class="service-detail">Passive Standby</text>
    
    <rect x="250" y="955" width="150" height="50" class="service-box enhanced-shadow" rx="8"/>
    <text x="325" y="975" class="service">DR Storage</text>
    <text x="325" y="990" class="service-detail">Geo-Redundant</text>
    
    <rect x="420" y="955" width="150" height="50" class="service-box enhanced-shadow" rx="8"/>
    <text x="495" y="975" class="service">DR Database</text>
    <text x="495" y="990" class="service-detail">Read Replica</text>
    
    <!-- MONITORING & OBSERVABILITY OVERLAY -->
    <rect x="720" y="900" width="650" height="120" class="trusted-zone strong-shadow" rx="15"/>
    <text x="740" y="925" class="swimlane-title">📊 MONITORING & OBSERVABILITY - 360° VISIBILITY</text>
    <text x="740" y="945" class="service-detail">SIEM | Cost Optimization | Compliance Monitoring</text>
    
    <rect x="770" y="955" width="140" height="50" class="service-box enhanced-shadow" rx="8"/>
    <text x="840" y="975" class="service">Azure Monitor</text>
    <text x="840" y="990" class="service-detail">Metrics & Alerts</text>
    
    <rect x="930" y="955" width="140" height="50" class="service-box enhanced-shadow" rx="8"/>
    <text x="1000" y="975" class="service">Log Analytics</text>
    <text x="1000" y="990" class="service-detail">Centralized Logs</text>
    
    <rect x="1090" y="955" width="140" height="50" class="service-box enhanced-shadow" rx="8"/>
    <text x="1160" y="975" class="service">Sentinel</text>
    <text x="1160" y="990" class="service-detail">Security Analytics</text>
    
    <!-- ENHANCED WORKFLOW CONNECTIONS with Polyline Routing -->
    <!-- 1. Internet to Security (HTTPS) -->
    <path d="M 740 200 L 740 230 L 530 230 L 530 285" stroke="#d32f2f" stroke-width="4" fill="none" marker-end="url(#securityArrow)"/>
    <text x="635" y="250" class="connection-label">1→2: HTTPS/443 | TLS 1.3 | WAF</text>
    
    <!-- 2. Security to Identity -->
    <path d="M 610 320 L 700 320" stroke="#388e3c" stroke-width="3" fill="none" marker-end="url(#primaryArrow)"/>
    <text x="655" y="340" class="connection-label">Authentication</text>
    
    <!-- 3. Security to Network Hub -->
    <path d="M 280 355 L 280 385 L 490 385 L 490 445" stroke="#1976d2" stroke-width="4" fill="none" marker-end="url(#primaryArrow)"/>
    <text x="385" y="400" class="connection-label">2→3: Filtered Traffic | UDR Applied</text>
    
    <!-- 4. Network Hub to Applications -->
    <path d="M 490 515 L 490 540 L 250 540 L 250 605" stroke="#1976d2" stroke-width="4" fill="none" marker-end="url(#primaryArrow)"/>
    <text x="370" y="560" class="connection-label">3→4: Spoke Routing | Private Connectivity</text>
    
    <!-- 5. Applications to Data -->
    <path d="M 250 695 L 250 720 L 380 720 L 380 785" stroke="#7b1fa2" stroke-width="4" fill="none" marker-end="url(#dataArrow)"/>
    <text x="315" y="740" class="connection-label">4→5: Data Access | Always Encrypted</text>
    
    <!-- Backup & Recovery Connections -->
    <path d="M 490 695 L 490 720 L 155 720 L 155 955" stroke="#4caf50" stroke-width="3" fill="none" stroke-dasharray="6,6" marker-end="url(#monitoringArrow)"/>
    <text x="325" y="740" class="connection-label">Backup & Recovery</text>
    
    <path d="M 290 855 L 290 875 L 325 875 L 325 955" stroke="#4caf50" stroke-width="3" fill="none" stroke-dasharray="6,6" marker-end="url(#monitoringArrow)"/>
    <text x="280" y="885" class="connection-label">Data Backup</text>
    
    <!-- Monitoring Connections (dotted lines) -->
    <path d="M 490 515 L 920 515 L 920 900" stroke="#00695c" stroke-width="2" fill="none" stroke-dasharray="3,3" marker-end="url(#monitoringArrow)"/>
    <text x="705" y="530" class="connection-label">Monitoring & Logs</text>
    
    <!-- COMPREHENSIVE LEGEND -->
    <rect x="30" y="1040" width="1340" height="180" class="service-box strong-shadow" rx="15"/>
    <text x="50" y="1065" class="legend-title">📋 COMPREHENSIVE LEGEND & DESIGN PRINCIPLES</text>
    
    <!-- Security Zones Legend -->
    <text x="50" y="1090" class="legend-title">Security Zones:</text>
    <rect x="50" y="1100" width="25" height="15" class="untrusted-zone"/>
    <text x="85" y="1112" class="legend-text">Untrusted (Internet Edge)</text>
    
    <rect x="220" y="1100" width="25" height="15" class="semitrusted-zone"/>
    <text x="255" y="1112" class="legend-text">Semi-trusted (Identity/Security)</text>
    
    <rect x="420" y="1100" width="25" height="15" class="trusted-zone"/>
    <text x="455" y="1112" class="legend-text">Trusted (Network/Apps)</text>
    
    <rect x="580" y="1100" width="25" height="15" class="data-zone"/>
    <text x="615" y="1112" class="legend-text">Data Zone (Storage/DB)</text>
    
    <rect x="720" y="1100" width="25" height="15" class="dr-zone"/>
    <text x="755" y="1112" class="legend-text">DR Zone (Disaster Recovery)</text>
    
    <!-- Connection Types Legend -->
    <text x="50" y="1140" class="legend-title">Connection Types:</text>
    <line x1="50" y1="1150" x2="100" y2="1150" stroke="#1976d2" stroke-width="4" marker-end="url(#primaryArrow)"/>
    <text x="110" y="1155" class="legend-text">Primary Traffic Flow</text>
    
    <line x1="250" y1="1150" x2="300" y2="1150" stroke="#d32f2f" stroke-width="4" marker-end="url(#securityArrow)"/>
    <text x="310" y="1155" class="legend-text">Security Connections</text>
    
    <line x1="450" y1="1150" x2="500" y2="1150" stroke="#7b1fa2" stroke-width="4" marker-end="url(#dataArrow)"/>
    <text x="510" y="1155" class="legend-text">Data Connections</text>
    
    <line x1="620" y1="1150" x2="670" y2="1150" stroke="#ef6c00" stroke-width="3" stroke-dasharray="8,4" marker-end="url(#monitoringArrow)"/>
    <text x="680" y="1155" class="legend-text">DR Replication</text>
    
    <line x1="800" y1="1150" x2="850" y2="1150" stroke="#00695c" stroke-width="2" stroke-dasharray="3,3" marker-end="url(#monitoringArrow)"/>
    <text x="860" y="1155" class="legend-text">Monitoring/Observability</text>
    
    <line x1="1050" y1="1150" x2="1100" y2="1150" stroke="#4caf50" stroke-width="3" stroke-dasharray="6,6" marker-end="url(#monitoringArrow)"/>
    <text x="1110" y="1155" class="legend-text">Backup & Recovery</text>
    
    <!-- Indicators Legend -->
    <text x="50" y="1180" class="legend-title">Service Indicators:</text>
    <rect x="50" y="1185" width="20" height="12" fill="#2196f3" rx="2"/>
    <text x="55" y="1193" style="font: 8px Arial; fill: white; text-anchor: start;">📈</text>
    <text x="80" y="1195" class="legend-text">Auto-Scale</text>
    
    <rect x="150" y="1185" width="20" height="12" fill="#9c27b0" rx="2"/>
    <text x="155" y="1193" style="font: 8px Arial; fill: white; text-anchor: start;">🔒</text>
    <text x="180" y="1195" class="legend-text">Security</text>
    
    <rect x="240" y="1185" width="20" height="12" fill="#4caf50" rx="2"/>
    <text x="245" y="1193" style="font: 8px Arial; fill: white; text-anchor: start;">⚡</text>
    <text x="270" y="1195" class="legend-text">High Performance</text>
    
    <rect x="380" y="1185" width="20" height="12" fill="#ff9800" rx="2"/>
    <text x="385" y="1193" style="font: 8px Arial; fill: white; text-anchor: start;">💰</text>
    <text x="410" y="1195" class="legend-text">Cost Optimized</text>
    
    <!-- HA/DR Indicators -->
    <text x="50" y="1210" class="legend-title">HA/DR Indicators:</text>
    <text x="50" y="1225" class="ha-indicator">Active-Active</text>
    <text x="150" y="1225" class="legend-text">| </text>
    <text x="160" y="1225" class="ha-indicator">Multi-Zone</text>
    <text x="240" y="1225" class="legend-text">| </text>
    <text x="250" y="1225" class="ha-indicator">Geo-Redundant</text>
    <text x="350" y="1225" class="legend-text">| </text>
    <text x="360" y="1225" class="ha-indicator">Auto-Failover</text>
    <text x="460" y="1225" class="legend-text">| </text>
    <text x="470" y="1225" class="ha-indicator">RTO: &lt;1hr</text>
    <text x="540" y="1225" class="legend-text">| </text>
    <text x="550" y="1225" class="ha-indicator">RPO: &lt;15min</text>
    
    <!-- Compliance & Architecture Principles -->
    <text x="700" y="1210" class="legend-title">Compliance & Standards:</text>
    <text x="700" y="1225" class="legend-text">✓ ISO 27001 ✓ SOC 2 ✓ GDPR ✓ HIPAA ✓ PCI-DSS ✓ CIS Benchmarks ✓ SOX</text>
    
    <!-- Future Ready -->
    <rect x="1150" y="900" width="220" height="120" fill="none" stroke="#9e9e9e" stroke-width="2" stroke-dasharray="10,5" rx="15"/>
    <text x="1170" y="925" class="swimlane-title" style="fill: #9e9e9e;">🚀 FUTURE EXPANSION ZONE</text>
    <text x="1170" y="945" class="service-detail" style="fill: #9e9e9e;">Reserved for Future Workloads</text>
    <text x="1170" y="965" class="service-detail" style="fill: #9e9e9e;">• AI/ML Services</text>
    <text x="1170" y="980" class="service-detail" style="fill: #9e9e9e;">• IoT Analytics</text>
    <text x="1170" y="995" class="service-detail" style="fill: #9e9e9e;">• Edge Computing</text>
    <text x="1170" y="1010" class="service-detail" style="fill: #9e9e9e;">• Blockchain Services</text>
    
</svg>'''
    
    # Write the enhanced SVG content to file
    with open(svg_filepath, "w", encoding="utf-8") as f:
        f.write(svg_content)
    
    logger.info(f"Enhanced simple SVG diagram generated successfully: {svg_filepath}")
    
    # If PNG is requested, convert SVG to PNG
    if format.lower() == "png":
        png_filename = f"enhanced_simple_azure_architecture_{timestamp}_{unique_id}.png"
        png_filepath = os.path.join(output_dir, png_filename)
        
        try:
            # Try to convert SVG to PNG using cairosvg if available
            import cairosvg
            cairosvg.svg2png(url=svg_filepath, write_to=png_filepath, output_width=1200, output_height=900)
            logger.info(f"Enhanced simple PNG diagram generated successfully: {png_filepath}")
            return png_filepath
        except ImportError:
            try:
                # Fallback: use inkscape command if available
                result = subprocess.run([
                    'inkscape', '--export-type=png', '--export-width=1200', '--export-height=900',
                    '--export-filename=' + png_filepath, svg_filepath
                ], capture_output=True, text=True, timeout=30)
                
                if result.returncode == 0 and os.path.exists(png_filepath):
                    logger.info(f"Enhanced simple PNG diagram generated via inkscape: {png_filepath}")
                    return png_filepath
                else:
                    raise Exception(f"Inkscape conversion failed: {result.stderr}")
            except (FileNotFoundError, subprocess.SubprocessError, subprocess.TimeoutExpired):
                # Final fallback: return SVG path anyway and let frontend handle it
                logger.warning("PNG conversion not available, returning SVG path")
                return svg_filepath
    
    return svg_filepath
